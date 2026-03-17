#!/usr/bin/env python3
"""
PPTX 2 PROFUNDO: Lección completa sobre cómo funciona la plataforma
Data Lake Governance Hub — con teoría, SQL real, ejemplos concretos
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== COLORES =====
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
SLATE_900 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_800 = RGBColor(0x1E, 0x40, 0x5F)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x05, 0x96, 0x69)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=SLATE_300, spacing=6):
    """Each line is (text, color, bold, size_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, clr, bld, sz = line, color, False, font_size
        else:
            txt = line[0]
            clr = line[1] if len(line) > 1 else color
            bld = line[2] if len(line) > 2 else False
            sz = line[3] if len(line) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = 'Calibri'
        p.space_after = Pt(spacing)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=SLATE_300):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.font.color.rgb = item[1] if len(item) > 1 else color
            p.font.bold = item[2] if len(item) > 2 else False
        else:
            p.text = item
            p.font.color.rgb = color
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def add_code_box(slide, left, top, width, height, code, font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p.font.name = 'Consolas'
    return shape

def add_badge(slide, left, top, text, bg_color, text_color=WHITE, width=1.2, height=0.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 2.0, 8.4, 1.5, title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 3.8, 8.4, 1.0, subtitle, font_size=16, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 6.5, 8.4, 0.5, "Gestión y Gobernanza de Datos — Universidad EAFIT — 2026", font_size=10, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(prs, section_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLATE_900)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.3), Inches(1.8), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, 0.5, 2.8, 9, 1.0, title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 0.5, 3.8, 9, 0.8, subtitle, font_size=14, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(prs, title, badge_text=None, badge_color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.5, 0.3, 7, 0.6, title, font_size=22, bold=True, color=WHITE)
    if badge_text:
        add_badge(slide, 8.0, 0.35, badge_text, badge_color, width=1.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SLATE_800
    shape.line.fill.background()
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # PORTADA
    # ================================================================
    title_slide(prs,
        "Data Lake Governance Hub\nLección Completa",
        "Cómo funciona la plataforma paso a paso\ncon teoría, SQL real y ejemplos concretos\n\ndatalake-governance-hub.vercel.app")

    # ================================================================
    # AGENDA
    # ================================================================
    slide = content_slide(prs, "Agenda de la lección")
    add_multiline(slide, 0.5, 1.1, 9, 6, [
        ("Módulo 1 — ¿Qué es un Data Lake y por qué necesita gobernanza?", BLUE, True, 15),
        ("    Concepto, Data Swamp, arquitectura Medallion", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 2 — Capa Bronze: los datos crudos", AMBER, True, 15),
        ("    JSONB, inmutabilidad, errores intencionales, cómo navegarla en la plataforma", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 3 — Capa Silver: limpieza y transformación", PURPLE, True, 15),
        ("    Casteo de tipos, COALESCE, INITCAP, quality_score, registros rechazados", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 4 — Capa Gold: datos listos para negocio", GREEN, True, 15),
        ("    Vistas materializadas, KPIs, dashboard, separación de capas", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 5 — Gobernanza: catálogo, calidad, linaje y auditoría", RED, True, 15),
        ("    Catálogo de datos, 6 dimensiones de calidad, trazabilidad, RBAC, PII", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 6 — Seguridad: RLS, enmascaramiento y roles", CYAN, True, 15),
        ("    Row Level Security, ventas_masked, 4 roles, principio de mínimo privilegio", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Módulo 7 — Operando la plataforma: guía práctica tab por tab", WHITE, True, 15),
        ("    Cómo usar cada sección, qué buscar, cómo interpretar los datos", SLATE_400, False, 12),
    ], spacing=2)

    # ================================================================
    # MÓDULO 1: QUÉ ES UN DATA LAKE
    # ================================================================
    section_slide(prs, 1, "¿Qué es un Data Lake?", "Y por qué sin gobernanza se convierte en un pantano")

    slide = content_slide(prs, "Definición de Data Lake", "CONCEPTO", BLUE)
    add_text_box(slide, 0.5, 1.1, 9, 0.8, "Un Data Lake es un repositorio centralizado que almacena TODOS los datos de una organización — estructurados, semi-estructurados y no estructurados — en su formato nativo y a cualquier escala.", font_size=15, color=WHITE)
    add_text_box(slide, 0.5, 2.2, 9, 0.4, "Diferencias clave con un Data Warehouse:", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 2.7, 4.3, 3, [
        ("Data Warehouse", AMBER, True),
        "Datos estructurados solamente",
        "Schema-on-write (defines antes de cargar)",
        "Optimizado para consultas SQL rápidas",
        "Costoso por GB almacenado",
        "Rígido — cambiar el esquema es difícil",
    ], font_size=12)
    add_bullet_list(slide, 5, 2.7, 4.5, 3, [
        ("Data Lake", BLUE, True),
        "Cualquier tipo de dato (JSON, CSV, imágenes, logs)",
        "Schema-on-read (defines al consultar)",
        "Flexible — almacena primero, analiza después",
        "Barato (object storage: ~$0.02/GB/mes)",
        "Riesgo: sin gobernanza = Data Swamp",
    ], font_size=12)
    add_text_box(slide, 0.5, 5.8, 9, 1, "Dato real: El 60% de los proyectos de Data Lake fracasan por falta de gobernanza (Gartner, 2023). Esta plataforma muestra cómo evitar ese fracaso.", font_size=12, color=SLATE_400)

    # Medallion Architecture
    slide = content_slide(prs, "Arquitectura Medallion (Bronze / Silver / Gold)", "ARQUITECTURA", BLUE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Popularizada por Databricks, es el estándar de la industria para organizar un Data Lake:", font_size=13, color=SLATE_300)

    # Bronze box
    add_badge(slide, 0.5, 1.8, "BRONZE", AMBER, width=2.5, height=0.4)
    add_multiline(slide, 0.5, 2.3, 2.5, 2.5, [
        ("Datos crudos", WHITE, True, 13),
        "Tal cual llegan del origen",
        "JSONB sin transformar",
        "Inmutables — nunca se modifican",
        "Particionados por fecha",
        "Incluyen errores y duplicados",
    ], font_size=11, spacing=3)

    # Arrow
    add_text_box(slide, 3.1, 2.8, 0.5, 0.5, "→", font_size=28, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    # Silver box
    add_badge(slide, 3.7, 1.8, "SILVER", PURPLE, width=2.5, height=0.4)
    add_multiline(slide, 3.7, 2.3, 2.5, 2.5, [
        ("Datos limpios", WHITE, True, 13),
        "Tipado correcto (DATE, INTEGER)",
        "Nulos tratados (COALESCE)",
        "Nombres normalizados (INITCAP)",
        "Duplicados eliminados",
        "Quality score calculado",
    ], font_size=11, spacing=3)

    # Arrow
    add_text_box(slide, 6.3, 2.8, 0.5, 0.5, "→", font_size=28, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    # Gold box
    add_badge(slide, 6.9, 1.8, "GOLD", EMERALD, width=2.5, height=0.4)
    add_multiline(slide, 6.9, 2.3, 2.5, 2.5, [
        ("Datos de negocio", WHITE, True, 13),
        "Vistas materializadas",
        "Agregaciones (SUM, AVG, COUNT)",
        "KPIs listos para dashboard",
        "Optimizado para consultas",
        "Accesible para analistas",
    ], font_size=11, spacing=3)

    add_text_box(slide, 0.5, 5.0, 9, 0.4, "En nuestra plataforma:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 5.5, 9, 1.5, [
        "Bronze: bronze.ventas_raw (8 registros JSONB) + bronze.sensores_raw (5 registros)",
        "Silver: silver.ventas_clean (7 registros tipados) + silver.sensores_clean (5 registros)",
        "Gold: gold.ventas_por_ciudad + gold.ventas_por_producto + gold.estado_sensores (vistas materializadas)",
    ], font_size=12)

    # ================================================================
    # MÓDULO 2: CAPA BRONZE
    # ================================================================
    section_slide(prs, 2, "Capa Bronze: los datos crudos", "Tab 'Bronze (Raw)' en la plataforma")

    slide = content_slide(prs, "¿Qué es Bronze y por qué es inmutable?", "BRONZE", AMBER)
    add_text_box(slide, 0.5, 1.1, 9, 0.8, "Bronze es la zona de aterrizaje (landing zone) del Data Lake. Los datos llegan TAL CUAL desde el sistema fuente — con errores, duplicados, formatos inconsistentes. NUNCA se modifican aquí.", font_size=14, color=SLATE_300)
    add_text_box(slide, 0.5, 2.1, 9, 0.4, "¿Por qué no modificar Bronze?", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 2.5, 9, 2, [
        "Si Silver tiene un bug, puedes reprocesar desde Bronze sin perder el dato original",
        "Evidencia legal: el dato original es prueba de lo que el sistema fuente envió",
        "Debugging: si un KPI parece raro, puedes comparar Bronze vs Silver para encontrar el problema",
        "Auditoría: reguladores pueden verificar que no se manipularon los datos",
    ], font_size=13)
    add_text_box(slide, 0.5, 4.5, 9, 0.4, "Estructura de la tabla bronze.ventas_raw:", font_size=14, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 4.9, 9, 2.2,
"""CREATE TABLE bronze.ventas_raw (
    id            BIGSERIAL PRIMARY KEY,    -- ID autoincremental
    raw_data      JSONB NOT NULL,           -- El dato CRUDO como JSON
    source_system TEXT NOT NULL,            -- "SAP-ERP", "iot-gateway"
    ingested_at   TIMESTAMPTZ DEFAULT now(),-- Cuándo se cargó
    ingested_by   TEXT,                     -- Quién/qué lo cargó ("etl-pipeline")
    file_name     TEXT,                     -- Archivo original ("export_ventas_20260301.json")
    batch_id      UUID DEFAULT gen_random_uuid() -- Agrupar cargas
);""", font_size=11)

    # Bronze - datos con errores
    slide = content_slide(prs, "Bronze en la plataforma: datos con errores reales", "ERRORES", RED)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Abre la tab 'Bronze (Raw)' — verás 8 registros JSONB. Busca estos 3 errores intencionales:", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Error 1 — Valor NULL (registro #4):", font_size=13, color=RED, bold=True)
    add_code_box(slide, 0.5, 2.2, 9, 0.8,
""""cantidad": null   ← No hay cantidad. ¿Cuántos se vendieron? Imposible saber.
La plataforma marca esto con etiqueta roja: "cantidad es NULL".""", font_size=11)

    add_text_box(slide, 0.5, 3.2, 9, 0.4, "Error 2 — String vacío (registro #4):", font_size=13, color=AMBER, bold=True)
    add_code_box(slide, 0.5, 3.6, 9, 0.8,
""""vendedor": ""   ← El vendedor está vacío. ¿Quién hizo la venta?
La plataforma marca esto con etiqueta amarilla: "vendedor vacío".""", font_size=11)

    add_text_box(slide, 0.5, 4.6, 9, 0.4, "Error 3 — Inconsistencia de formato (registro #4):", font_size=13, color=BLUE, bold=True)
    add_code_box(slide, 0.5, 5.0, 9, 0.8,
""""ciudad": "medellin"   ← minúsculas, sin tilde.
Los otros registros dicen "Medellín". Esto rompe GROUP BY y filtros.""", font_size=11)

    add_text_box(slide, 0.5, 6.1, 9, 0.8, "Nota: En la vida real estos errores son COMUNES. Los sistemas ERP, sensores IoT y APIs externas envían datos sucios todo el tiempo. El Data Lake debe saber manejarlos.", font_size=12, color=SLATE_400)

    # ================================================================
    # MÓDULO 3: CAPA SILVER
    # ================================================================
    section_slide(prs, 3, "Capa Silver: limpieza y transformación", "Tab 'Silver (Clean)' y tab 'Bronze vs Silver'")

    slide = content_slide(prs, "¿Qué pasa en Bronze → Silver?", "SILVER", PURPLE)
    add_text_box(slide, 0.5, 1.1, 9, 0.4, "Silver es donde los datos se limpian, validan y tipan. Es la capa más importante del pipeline.", font_size=14, color=SLATE_300)
    add_text_box(slide, 0.5, 1.7, 9, 0.4, "Las 5 transformaciones que aplicamos:", font_size=14, color=PURPLE, bold=True)

    add_text_box(slide, 0.5, 2.2, 4.3, 0.3, "1. Casteo de tipos", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 0.5, 2.5, 4.3, 0.6, "(raw_data->>'fecha')::date\n(raw_data->>'cantidad')::integer", font_size=10)

    add_text_box(slide, 5, 2.2, 4.5, 0.3, "2. Tratamiento de vacíos", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 5, 2.5, 4.5, 0.6, "COALESCE(NULLIF(vendedor, ''),\n  'Sin asignar')", font_size=10)

    add_text_box(slide, 0.5, 3.3, 4.3, 0.3, "3. Normalización de texto", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 0.5, 3.6, 4.3, 0.6, 'INITCAP(raw_data->>\'ciudad\')\n"medellin" --> "Medellin"', font_size=10)

    add_text_box(slide, 5, 3.3, 4.5, 0.3, "4. Filtrado de inválidos", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 5, 3.6, 4.5, 0.6, "WHERE cantidad IS NOT NULL\n  AND cantidad > 0", font_size=10)

    add_text_box(slide, 0.5, 4.4, 4.3, 0.3, "5. Quality Score", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 0.5, 4.7, 4.3, 0.6, "CASE WHEN vendedor = ''\n  THEN 0.70 ELSE 1.00 END", font_size=10)

    add_text_box(slide, 5, 4.4, 4.5, 0.3, "6. Campo calculado", font_size=13, color=WHITE, bold=True)
    add_code_box(slide, 5, 4.7, 4.5, 0.6, "ingreso_total GENERATED ALWAYS AS\n  (cantidad * precio_unit) STORED", font_size=10)

    add_text_box(slide, 0.5, 5.6, 9, 0.4, "Resultado:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 6.0, 9, 1, [
        "8 registros Bronze → 7 aprobados en Silver + 1 rechazado (registro #4: cantidad NULL)",
        "El registro rechazado PERMANECE en Bronze — nunca se borra. Silver solo tiene datos válidos.",
    ], font_size=12)

    # Compare tab explanation
    slide = content_slide(prs, "Tab 'Bronze vs Silver': cómo interpretar la comparación", "COMPARAR", CYAN)
    add_text_box(slide, 0.5, 1.1, 9, 0.8, "Esta tab muestra cada registro Bronze lado a lado con su versión Silver. Es la forma más visual de entender qué cambió y por qué.", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 2.1, 9, 0.4, "Cómo leer cada tarjeta:", font_size=14, color=CYAN, bold=True)
    add_bullet_list(slide, 0.5, 2.5, 9, 3.5, [
        ("Encabezado: B#1 → S#1 = Bronze ID 1 se transformó en Silver ID 1", WHITE, False),
        ("Etiqueta verde 'LIMPIO' = el registro pasó sin problemas", GREEN, False),
        ("Etiqueta amarilla '1 transformación(es)' = se corrigió algo (ej: vendedor vacío)", AMBER, False),
        ("Etiqueta roja 'RECHAZADO' = el registro NO pasó a Silver", RED, False),
        "Lado izquierdo: datos Bronze con campos problemáticos en ROJO",
        "Lado derecho: datos Silver con campos corregidos en VERDE",
        "Footer: detalle específico de cada transformación aplicada",
    ], font_size=13)

    add_text_box(slide, 0.5, 5.8, 9, 0.4, "Ejercicio para los estudiantes:", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 6.2, 9, 1, [
        "Abre la tab 'Bronze vs Silver'. Encuentra el registro RECHAZADO. ¿Por qué se rechazó? ¿Qué 3 problemas tenía?",
    ], font_size=13)

    # Silver tab details
    slide = content_slide(prs, "Tab 'Silver (Clean)': los datos ya validados", "SILVER", PURPLE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Aquí ves la tabla silver.ventas_clean con todas las columnas tipadas:", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.7, 9, 0.4, "Columnas importantes:", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 2.1, 9, 3.5, [
        "fecha (DATE) — casteo correcto, permite filtrar por rango de fechas",
        "cantidad (INTEGER, CHECK > 0) — validación a nivel de base de datos, no solo aplicación",
        "precio_unit (NUMERIC(12,2), CHECK > 0) — precisión decimal para finanzas",
        "vendedor — 'Sin asignar' en amarillo indica que el dato original estaba vacío",
        "ciudad — normalizada con INITCAP(): consistente para GROUP BY",
        "ingreso_total — campo GENERATED ALWAYS: se calcula solo, no puede tener errores manuales",
        "quality_score — 1.00 = perfecto, 0.70 = dato con algún problema corregido",
        "source_record_id — enlace directo al registro Bronze original (trazabilidad/linaje)",
    ], font_size=12)

    add_text_box(slide, 0.5, 5.5, 9, 0.4, "Contadores en la parte superior:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 5.9, 9, 1.5, [
        "8 registros Bronze (entrada) → 7 registros Silver (aprobados) → 1 rechazado por calidad",
        "Esto demuestra que el pipeline de calidad funciona: no todo lo que llega pasa automáticamente",
    ], font_size=12)

    # ================================================================
    # MÓDULO 4: CAPA GOLD
    # ================================================================
    section_slide(prs, 4, "Capa Gold: datos listos para negocio", "Tab 'Dashboard Gold'")

    slide = content_slide(prs, "¿Qué son las vistas materializadas de Gold?", "GOLD", GREEN)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Gold no almacena datos nuevos — agrega y resume los datos de Silver para consumo de negocio.", font_size=14, color=SLATE_300)
    add_text_box(slide, 0.5, 1.7, 9, 0.3, "Vista materializada = una consulta SQL pre-calculada que se guarda como tabla:", font_size=13, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 2.1, 9, 2.5,
"""-- Esta consulta se ejecutó UNA VEZ y el resultado se guardó como tabla
CREATE MATERIALIZED VIEW gold.ventas_por_ciudad AS
SELECT
    ciudad,
    COUNT(*)           AS num_transacciones,  -- ¿Cuántas ventas?
    SUM(cantidad)      AS unidades_vendidas,  -- ¿Cuántas unidades?
    SUM(ingreso_total) AS ingreso_total,      -- ¿Cuánto dinero?
    ROUND(AVG(ingreso_total), 2) AS ticket_promedio,   -- Promedio por venta
    ROUND(AVG(quality_score), 2) AS calidad_promedio,   -- Calidad de los datos
    MIN(fecha) AS primera_venta,
    MAX(fecha) AS ultima_venta
FROM silver.ventas_clean    -- Lee de Silver, NUNCA de Bronze
GROUP BY ciudad
ORDER BY ingreso_total DESC;""", font_size=11)

    add_text_box(slide, 0.5, 4.8, 9, 0.4, "¿Por qué vistas MATERIALIZADAS y no vistas normales?", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 5.2, 9, 2, [
        "Vista normal: ejecuta el GROUP BY cada vez que alguien consulta — lento con millones de filas",
        "Vista materializada: el resultado ya está calculado — consultas instantáneas",
        "Se refresca con: REFRESH MATERIALIZED VIEW gold.ventas_por_ciudad;",
        "Trade-off: los datos Gold no se actualizan en tiempo real, sino cuando se refresca",
    ], font_size=12)

    # Dashboard tab
    slide = content_slide(prs, "Tab 'Dashboard Gold': cómo operarlo", "DASHBOARD", GREEN)
    add_text_box(slide, 0.5, 1.1, 9, 0.4, "Al abrir la plataforma, esta es la primera tab. Muestra 4 KPIs y 2 tablas:", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.7, 4.3, 0.3, "KPI 1: Ingreso Total", font_size=13, color=GREEN, bold=True)
    add_text_box(slide, 0.5, 2.0, 4.3, 0.5, "Suma de (cantidad × precio) de todas las ventas aprobadas. En COP.", font_size=11, color=SLATE_400)

    add_text_box(slide, 5, 1.7, 4.5, 0.3, "KPI 2: Transacciones", font_size=13, color=BLUE, bold=True)
    add_text_box(slide, 5, 2.0, 4.5, 0.5, "Cantidad total de ventas que pasaron de Bronze a Silver.", font_size=11, color=SLATE_400)

    add_text_box(slide, 0.5, 2.6, 4.3, 0.3, "KPI 3: Calidad Promedio", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 0.5, 2.9, 4.3, 0.5, "Promedio ponderado de quality_score. Si baja de 80% = problema.", font_size=11, color=SLATE_400)

    add_text_box(slide, 5, 2.6, 4.5, 0.3, "KPI 4: Registros Rechazados", font_size=13, color=RED, bold=True)
    add_text_box(slide, 5, 2.9, 4.5, 0.5, "Diferencia entre Bronze y Silver. Si crece = problema en la fuente.", font_size=11, color=SLATE_400)

    add_text_box(slide, 0.5, 3.7, 9, 0.3, "Tabla 'Ventas por Ciudad':", font_size=14, color=WHITE, bold=True)
    add_bullet_list(slide, 0.5, 4.1, 9, 1.5, [
        "Columna 'Calidad' con badge verde (>=90%) o amarillo (<90%)",
        "Si Medellín muestra 85%, es porque uno de sus registros tenía vendedor vacío → score 0.70",
        "Esto permite identificar ciudades con datos de menor confiabilidad",
    ], font_size=12)

    add_text_box(slide, 0.5, 5.5, 9, 0.3, "Regla de oro del dashboard:", font_size=14, color=RED, bold=True)
    add_text_box(slide, 0.5, 5.9, 9, 0.8, "El código del dashboard dice: supabase.from('gold_ventas_por_ciudad').select('*')\nNUNCA dice: supabase.from('bronze_ventas_raw')\n\nUn analista de negocio no debería tocar Bronze ni Silver. Solo Gold.", font_size=12, color=SLATE_300)

    # ================================================================
    # MÓDULO 5: GOBERNANZA
    # ================================================================
    section_slide(prs, 5, "Gobernanza del Data Lake", "Lo que separa un proyecto exitoso de un fracaso")

    # Catálogo
    slide = content_slide(prs, "Tab 'Catálogo': el registro de todos los datasets", "CATÁLOGO", BLUE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Regla #1 de Gobernanza: Ningún dato entra al Lake sin metadatos. El catálogo es la 'guía telefónica' de tus datos.", font_size=14, color=SLATE_300)
    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Para cada dataset, el catálogo registra:", font_size=14, color=BLUE, bold=True)
    add_bullet_list(slide, 0.5, 2.2, 4.3, 3.5, [
        ("Esquema + tabla", WHITE, True),
        "bronze.ventas_raw, silver.ventas_clean, etc.",
        "",
        ("Clasificación de sensibilidad", WHITE, True),
        "public: cualquiera puede ver",
        "internal: solo empleados",
        "confidential: datos financieros/PII",
        "restricted: datos regulados",
    ], font_size=11)
    add_bullet_list(slide, 5, 2.2, 4.5, 3.5, [
        ("PII (Personally Identifiable Information)", WHITE, True),
        "ventas_raw: PII=Sí (tiene NIT, nombre)",
        "estado_sensores: PII=No (solo máquinas)",
        "",
        ("Freshness (frescura)", WHITE, True),
        "daily: se actualiza cada día",
        "real-time: datos en vivo",
        "hourly: cada hora",
    ], font_size=11)
    add_text_box(slide, 0.5, 5.8, 9, 0.4, "Cómo leer la barra de calidad:", font_size=13, color=GREEN, bold=True)
    add_text_box(slide, 0.5, 6.2, 9, 0.6, "Cada dataset tiene una barra visual: verde (>=95%), azul (>=90%), ámbar (>=85%), rojo (<85%). Bronze.ventas_raw tiene 85% porque incluye datos crudos con errores. Gold.ventas_por_ciudad tiene 98% porque ya pasó por limpieza.", font_size=11, color=SLATE_400)

    # Calidad
    slide = content_slide(prs, "Tab 'Calidad': las 6 dimensiones de DAMA-DMBOK2", "CALIDAD", AMBER)
    add_text_box(slide, 0.5, 1.1, 9, 0.4, "El framework DAMA-DMBOK2 define 6 dimensiones de calidad de datos:", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.7, 4.3, 0.3, "1. Completitud", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 0.5, 2.0, 4.3, 0.5, "¿Faltan valores? Regla: cantidad NOT NULL.\nEn la plataforma: 1 registro rechazado por NULL.", font_size=10, color=SLATE_400)

    add_text_box(slide, 5, 1.7, 4.5, 0.3, "2. Validez", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 5, 2.0, 4.5, 0.5, "¿Los valores son posibles? Regla: cantidad > 0.\nUna venta de -5 unidades no tiene sentido.", font_size=10, color=SLATE_400)

    add_text_box(slide, 0.5, 2.7, 4.3, 0.3, "3. Consistencia", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 0.5, 3.0, 4.3, 0.5, "¿Mismo formato siempre? Regla: INITCAP(ciudad).\n'medellin', 'Medellín', 'MEDELLIN' → todos iguales.", font_size=10, color=SLATE_400)

    add_text_box(slide, 5, 2.7, 4.5, 0.3, "4. Unicidad", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 5, 3.0, 4.5, 0.5, "¿Hay duplicados? Regla: deduplicación por ID.\nEvita contar la misma venta dos veces.", font_size=10, color=SLATE_400)

    add_text_box(slide, 0.5, 3.7, 4.3, 0.3, "5. Oportunidad", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 0.5, 4.0, 4.3, 0.5, "¿Los datos están frescos? En catálogo: freshness.\nSi Gold muestra datos de enero en marzo = problema.", font_size=10, color=SLATE_400)

    add_text_box(slide, 5, 3.7, 4.5, 0.3, "6. Confiabilidad", font_size=13, color=AMBER, bold=True)
    add_text_box(slide, 5, 4.0, 4.5, 0.5, "¿Puedo confiar? En la plataforma: quality_score.\nSi < 0.80 = alerta, el dato es sospechoso.", font_size=10, color=SLATE_400)

    add_text_box(slide, 0.5, 4.8, 9, 0.4, "En la tab Calidad ves estas 3 secciones:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 5.2, 9, 2, [
        "Arriba: 5 métricas (Bronze count, Silver count, Rechazados, Tasa aprobación, Quality Score promedio)",
        "Centro-izquierda: 7 reglas de calidad activas con tipo (Completitud/Validez/Consistencia/etc.)",
        "Centro-derecha: Problemas detectados con severidad (Crítico=rojo, Warning=amarillo, Info=azul)",
        "Abajo: Barras de quality score por dataset del catálogo",
    ], font_size=12)

    # Linaje
    slide = content_slide(prs, "Tab 'Linaje': de dónde viene cada dato", "LINAJE", PURPLE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "El linaje responde: 'Este KPI de $22.250.000 en Barranquilla... ¿de dónde sale? ¿Puedo confiar?'", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Parte 1 — Grafo de linaje (flujo macro):", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 2.2, 9, 2, [
        "FUENTE → BRONZE → SILVER → GOLD → CONSUMO, con cada tabla nombrada",
        "Puedes ver qué transformaciones ocurren entre cada capa",
        "En producción esto lo genera automáticamente Apache Atlas, OpenLineage o dbt lineage",
    ], font_size=12)

    add_text_box(slide, 0.5, 3.8, 9, 0.4, "Parte 2 — Linaje a nivel de registro:", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 4.2, 9, 1.5, [
        "Tabla que muestra: Silver ID → Bronze ID → Sistema fuente → Archivo original",
        "Ejemplo: Silver #1 (Sensor IoT X200, Medellín, calidad 100%) ← Bronze #1 (SAP-ERP, export_ventas_20260301.json)",
    ], font_size=12)

    add_text_box(slide, 0.5, 5.5, 9, 0.4, "¿Cómo funciona técnicamente?", font_size=14, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 5.9, 9, 1.2,
"""-- La columna source_record_id en Silver apunta a Bronze:
source_record_id BIGINT REFERENCES bronze.ventas_raw(id)

-- Para rastrear un dato Gold → Silver → Bronze:
-- Gold: "Barranquilla vendió $22.250.000"
-- → Silver: SELECT * FROM silver.ventas_clean WHERE ciudad = 'Barranquilla'
-- → Bronze: SELECT * FROM bronze.ventas_raw WHERE id = {source_record_id}""", font_size=10)

    # ================================================================
    # MÓDULO 6: SEGURIDAD
    # ================================================================
    section_slide(prs, 6, "Seguridad: RLS, enmascaramiento y roles", "Tabs 'Auditoría' y 'Usuarios'")

    # Auditoría
    slide = content_slide(prs, "Tab 'Auditoría': quién accedió a qué y cuándo", "AUDITORÍA", RED)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "La auditoría no es opcional — es un requisito legal (Ley 1581 Habeas Data en Colombia, GDPR en Europa, SOX para financieros).", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Cada fila del log registra:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 2.2, 4.3, 3, [
        ("Timestamp", WHITE, True),
        "Fecha y hora exacta del acceso",
        "",
        ("Usuario + Rol", WHITE, True),
        "Email y rol (data_engineer, analyst...)",
        "",
        ("Acción", WHITE, True),
        "SELECT (consultó), INSERT (cargó), EXPORT (descargó)",
    ], font_size=11)
    add_bullet_list(slide, 5, 2.2, 4.5, 3, [
        ("Tabla accedida", WHITE, True),
        "Esquema + nombre (gold.ventas_por_ciudad)",
        "",
        ("Filas afectadas", WHITE, True),
        "Cuántos registros consultó o cargó",
        "",
        ("IP de origen", WHITE, True),
        "Desde qué red se conectó",
    ], font_size=11)

    add_text_box(slide, 0.5, 5.0, 9, 0.4, "Escenario real: ¿para qué sirve?", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 5.4, 9, 1.8, [
        "Caso 1: Un analista descargó (EXPORT) datos confidenciales. ¿Cuáles? ¿Cuántos registros? ¿Desde dónde?",
        "Caso 2: ¿Quién cargó los datos de marzo? → Filtrar por action='INSERT', ver ingested_by='etl-pipeline'",
        "Caso 3: Un regulador pide evidencia de que solo personal autorizado accedió a datos PII",
    ], font_size=12)

    # RBAC y usuarios
    slide = content_slide(prs, "Tab 'Usuarios': RBAC y principio de mínimo privilegio", "RBAC", CYAN)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "RBAC = Role-Based Access Control. Cada usuario tiene el mínimo acceso necesario para hacer su trabajo.", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Los 4 roles y lo que pueden hacer:", font_size=14, color=CYAN, bold=True)

    add_text_box(slide, 0.5, 2.3, 4.3, 0.3, "ADMIN (Admin Lake)", font_size=13, color=RED, bold=True)
    add_text_box(slide, 0.5, 2.6, 4.3, 0.4, "Todo: Bronze + Silver + Gold + Audit + Admin.\nSolo el equipo de TI/gobernanza.", font_size=10, color=SLATE_400)

    add_text_box(slide, 5, 2.3, 4.5, 0.3, "DATA ENGINEER (Carlos Martínez)", font_size=13, color=GREEN, bold=True)
    add_text_box(slide, 5, 2.6, 4.5, 0.4, "Bronze R/W + Silver R/W + Gold R.\nCarga datos y construye pipelines.", font_size=10, color=SLATE_400)

    add_text_box(slide, 0.5, 3.2, 4.3, 0.3, "DATA SCIENTIST (María López)", font_size=13, color=PURPLE, bold=True)
    add_text_box(slide, 0.5, 3.5, 4.3, 0.4, "Silver MASKED (R) + Gold (R) + Sandbox.\nExplora datos SIN ver PII real.", font_size=10, color=SLATE_400)

    add_text_box(slide, 5, 3.2, 4.5, 0.3, "ANALYST (Pedro García)", font_size=13, color=BLUE, bold=True)
    add_text_box(slide, 5, 3.5, 4.5, 0.4, "Gold (R) solamente.\nSolo ve datos agregados y limpios.", font_size=10, color=SLATE_400)

    add_text_box(slide, 0.5, 4.2, 9, 0.4, "Enmascaramiento de PII (datos personales):", font_size=14, color=RED, bold=True)
    add_code_box(slide, 0.5, 4.6, 9, 1.8,
"""-- Vista silver.ventas_masked — lo que ve un Data Scientist:
CREATE VIEW silver.ventas_masked AS
SELECT
    id, fecha, producto, cantidad, precio_unit, vendedor, ciudad,
    CONCAT(LEFT(cliente_nit, 3), '****', RIGHT(cliente_nit, 2))
        AS cliente_nit_masked,       -- "900123456" → "900****56"
    CONCAT(LEFT(cliente_nombre, 3), '***')
        AS cliente_nombre_masked,    -- "Acueducto Municipal" → "Acu***"
    ingreso_total, quality_score
FROM silver.ventas_clean;
-- María (data_scientist) ve esta vista, NO la tabla completa""", font_size=10)

    # RLS
    slide = content_slide(prs, "Row Level Security (RLS): control a nivel de fila", "RLS", RED)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "RLS es una feature nativa de PostgreSQL (y Supabase) que controla qué FILAS puede ver cada usuario.", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "¿Cómo funciona?", font_size=14, color=RED, bold=True)
    add_code_box(slide, 0.5, 2.2, 9, 2,
"""-- 1. Habilitar RLS en la tabla
ALTER TABLE silver.ventas_clean ENABLE ROW LEVEL SECURITY;

-- 2. Crear política: quién puede ver qué
CREATE POLICY "engineers_full_access"
    ON silver.ventas_clean
    FOR ALL                     -- SELECT, INSERT, UPDATE, DELETE
    USING (
        EXISTS (
            SELECT 1 FROM public.lake_users
            WHERE email = current_setting('app.current_user_email')
            AND role IN ('data_engineer', 'admin')
        )
    );
-- Solo data_engineers y admins pueden ver/editar Silver""", font_size=10)

    add_text_box(slide, 0.5, 4.5, 9, 0.4, "Sin RLS vs Con RLS:", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 4.9, 4.3, 2, [
        ("Sin RLS:", RED, True),
        "Todos ven todo",
        "Un becario puede borrar Bronze",
        "Un analista puede leer NITs reales",
        "No hay control granular",
    ], font_size=11)
    add_bullet_list(slide, 5, 4.9, 4.5, 2, [
        ("Con RLS:", GREEN, True),
        "Cada rol ve solo lo permitido",
        "Bronze es intocable para no-engineers",
        "PII enmascarado para data scientists",
        "Auditable y demostrable",
    ], font_size=11)

    # ================================================================
    # MÓDULO 7: OPERANDO LA PLATAFORMA
    # ================================================================
    section_slide(prs, 7, "Operando la plataforma", "Guía práctica: qué hacer en cada tab")

    # Tab Sensores
    slide = content_slide(prs, "Tab 'Sensores IoT': monitoreo con detección de anomalías", "SENSORES", GREEN)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Esta tab consume gold.estado_sensores — una vista materializada que agrega las lecturas de cada sensor.", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 9, 0.4, "Cada tarjeta de sensor muestra:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 2.2, 9, 2, [
        "ID del sensor (SEN-001, SEN-002...) y ubicación (Planta Norte, Red Principal...)",
        "Valor promedio, mínimo y máximo de la medición + unidad (L/s, bar, pH, m)",
        "Batería: verde (>70%), ámbar (50-70%), rojo (<50%) — SEN-004 tiene 45% = alerta",
        "Indicador de anomalía: badge rojo 'Anomalía' si alguna lectura es físicamente imposible",
    ], font_size=12)

    add_text_box(slide, 0.5, 4.0, 9, 0.4, "Ejemplo de anomalía — SEN-001:", font_size=14, color=RED, bold=True)
    add_code_box(slide, 0.5, 4.4, 9, 1.8,
"""-- En Bronze, SEN-001 tiene dos lecturas:
-- Lectura 1: valor = 245.7 L/s (normal)
-- Lectura 2: valor = -5.0 L/s  (IMPOSIBLE — caudal negativo)

-- El pipeline Bronze → Silver marcó:
CASE WHEN (raw_data->>'valor')::numeric < 0
     THEN TRUE ELSE FALSE END AS is_anomaly

-- En Gold, la vista materializada agrega:
BOOL_OR(is_anomaly) AS tiene_anomalias  -- TRUE si CUALQUIER lectura es anómala
-- Resultado: SEN-001 muestra badge "Anomalía detectada" en la plataforma""", font_size=10)

    add_text_box(slide, 0.5, 6.5, 9, 0.5, "En producción: estas anomalías dispararían alertas por Slack/email al equipo de operaciones.", font_size=12, color=SLATE_400)

    # Tab Arquitectura
    slide = content_slide(prs, "Tab 'Arquitectura': vista técnica del sistema", "STACK", PURPLE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Esta tab resume todo el sistema: el flujo Medallion, la seguridad implementada, la gobernanza y el stack.", font_size=14, color=SLATE_300)

    add_text_box(slide, 0.5, 1.8, 4.3, 0.3, "Flujo Medallion visual:", font_size=13, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 2.1, 4.3, 1.5, [
        "3 bloques: Bronze → Silver → Gold",
        "Cada bloque lista sus tablas específicas",
        "Muestra conteo de registros en tiempo real",
    ], font_size=11)

    add_text_box(slide, 5, 1.8, 4.5, 0.3, "Checklist de seguridad:", font_size=13, color=RED, bold=True)
    add_bullet_list(slide, 5, 2.1, 4.5, 1.5, [
        "RLS, Enmascaramiento PII, RBAC 4 roles",
        "Auditoría con IP, Cifrado TLS + AES-256",
        "Todas marcadas con check verde",
    ], font_size=11)

    add_text_box(slide, 0.5, 3.8, 4.3, 0.3, "Checklist de gobernanza:", font_size=13, color=BLUE, bold=True)
    add_bullet_list(slide, 0.5, 4.1, 4.3, 1.5, [
        "Catálogo con metadatos obligatorios",
        "Clasificación, Quality Score, Linaje",
        "Data Stewardship (owners por dataset)",
    ], font_size=11)

    add_text_box(slide, 5, 3.8, 4.5, 0.3, "Stack tecnológico:", font_size=13, color=GREEN, bold=True)
    add_bullet_list(slide, 5, 4.1, 4.5, 1.5, [
        "Supabase: PostgreSQL + RLS + Auth",
        "Next.js 14: React SSR framework",
        "Vercel: Edge deployment",
        "Tailwind CSS: Utility-first styling",
    ], font_size=11)

    add_text_box(slide, 0.5, 5.8, 9, 0.4, "¿Por qué este stack para un ejercicio de Gobernanza?", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 6.2, 9, 1, [
        "Supabase = PostgreSQL real con RLS nativo (no un mock) — misma tecnología que usan empresas como GitHub, Twitch",
        "Vercel = deploy en 1 clic, HTTPS automático, CDN global — así funciona producción real",
    ], font_size=12)

    # Guía de navegación
    slide = content_slide(prs, "Flujo de trabajo recomendado para explorar", "GUÍA", BLUE)
    add_text_box(slide, 0.5, 1.1, 9, 0.5, "Sigue este orden para entender el sistema completo:", font_size=14, color=SLATE_300)

    add_multiline(slide, 0.5, 1.8, 9, 5.5, [
        ("Paso 1 → Tab 'Bronze'", AMBER, True, 15),
        ("    Ve los datos crudos. Identifica los 3 errores (NULL, vacío, minúsculas).", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 2 → Tab 'Bronze vs Silver'", CYAN, True, 15),
        ("    Compara antes/después. Encuentra el registro rechazado. Entiende cada transformación.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 3 → Tab 'Silver'", PURPLE, True, 15),
        ("    Ve los datos limpios. Nota los tipos de datos, el quality_score y el campo calculado.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 4 → Tab 'Dashboard Gold'", GREEN, True, 15),
        ("    Ve los KPIs. ¿El ingreso total tiene sentido? ¿Por qué Medellín tiene 85% de calidad?", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 5 → Tab 'Calidad'", AMBER, True, 15),
        ("    Revisa las reglas activas. ¿Qué problemas se detectaron? ¿Cuál es la tasa de aprobación?", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 6 → Tab 'Linaje'", PURPLE, True, 15),
        ("    Rastrea un dato Gold hasta Bronze. ¿De qué archivo vino? ¿Qué sistema lo generó?", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("Paso 7 → Tabs 'Catálogo', 'Auditoría', 'Usuarios', 'Sensores', 'Arquitectura'", WHITE, True, 15),
        ("    Explora la gobernanza completa: metadatos, accesos, roles, anomalías, stack técnico.", SLATE_400, False, 12),
    ], spacing=1)

    # Preguntas
    slide = content_slide(prs, "Preguntas de comprensión", "REFLEXIÓN", AMBER)
    add_multiline(slide, 0.5, 1.1, 9, 6, [
        ("1. ¿Por qué el registro #4 fue rechazado?", WHITE, True, 14),
        ("    Porque tenía cantidad = NULL. La regla de Silver dice: cantidad IS NOT NULL AND > 0.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("2. ¿Qué hace INITCAP() y por qué importa?", WHITE, True, 14),
        ('    Convierte "medellin" en "Medellin". Sin esto, GROUP BY los trata como ciudades diferentes.', SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("3. ¿Por qué Medellín tiene quality_score de 85% y no 100%?", WHITE, True, 14),
        ("    Porque uno de sus registros tenía vendedor vacío → score 0.70 → promedio baja.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("4. ¿Qué ve María (data_scientist) vs Pedro (analyst)?", WHITE, True, 14),
        ("    María ve Silver con PII enmascarado + Gold. Pedro solo ve Gold.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("5. Si el ERP agrega un campo 'descuento', ¿se rompe Bronze?", WHITE, True, 14),
        ("    No. Bronze usa JSONB, que acepta campos nuevos sin cambiar la tabla.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("6. ¿Para qué sirve source_record_id en Silver?", WHITE, True, 14),
        ("    Para linaje: rastrear cada dato limpio hasta su registro crudo original en Bronze.", SLATE_400, False, 12),
        ("", WHITE, False, 6),
        ("7. ¿Cuándo necesitas refrescar una vista materializada?", WHITE, True, 14),
        ("    Cuando Silver cambia. Comando: REFRESH MATERIALIZED VIEW gold.ventas_por_ciudad;", SLATE_400, False, 12),
    ], spacing=1)

    # Conceptos clave resumen
    slide = content_slide(prs, "Conceptos clave — Resumen de la lección", "RESUMEN", GREEN)
    add_multiline(slide, 0.5, 1.1, 9, 6.2, [
        ("Arquitectura Medallion", WHITE, True, 15),
        ("Bronze (crudo) → Silver (limpio) → Gold (negocio). Separación de responsabilidades.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Inmutabilidad en Bronze", WHITE, True, 15),
        ("Nunca modificar datos crudos. Si hay errores, corregirlos en Silver. Bronze es evidencia.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Schema-on-read vs Schema-on-write", WHITE, True, 15),
        ("Bronze acepta cualquier JSON. Silver impone esquema estricto con tipos y validaciones.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Quality Score", WHITE, True, 15),
        ("Métrica numérica por registro que indica confiabilidad. Permite alertar sin bloquear.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Data Lineage (Linaje)", WHITE, True, 15),
        ("Rastrear cada dato hasta su origen. Implementado con source_record_id en Silver.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Row Level Security (RLS)", WHITE, True, 15),
        ("Control de acceso a nivel de FILA en PostgreSQL. Cada rol ve solo lo que debe ver.", SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("PII Masking (Enmascaramiento)", WHITE, True, 15),
        ('Vista que oculta datos personales: NIT "900123456" → "900****56". Para data scientists.', SLATE_400, False, 12),
        ("", WHITE, False, 4),
        ("Catálogo de Datos", WHITE, True, 15),
        ("Registro obligatorio de cada dataset: dueño, clasificación, PII, calidad, frescura.", SLATE_400, False, 12),
    ], spacing=1)

    # Cierre
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLATE_900)
    add_text_box(slide, 0.8, 1.8, 8.4, 1, "Explora la plataforma", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 3.0, 8.4, 1, "datalake-governance-hub.vercel.app", font_size=22, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 4.0, 8.4, 1.5, "Navega por las 11 tabs siguiendo el flujo recomendado.\nCompara Bronze vs Silver, rastrea el linaje,\nrevisa la auditoría, entiende los roles.\n\nRecuerda: un Data Lake sin gobernanza\nes solo un bucket con datos.", font_size=15, color=SLATE_300, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 6.2, 8.4, 1, "Gestión y Gobernanza de Datos — Universidad EAFIT — 2026", font_size=12, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    # Save
    path = os.path.expanduser("~/Documents/datalake-governance-hub/PPTX_2_Como_Funciona_Plataforma.pptx")
    prs.save(path)
    print(f"✅ PPTX 2 guardado: {path}")
    return path


if __name__ == '__main__':
    path = build()
    print(f"\n🎉 PowerPoint profundo creado: {path}")
