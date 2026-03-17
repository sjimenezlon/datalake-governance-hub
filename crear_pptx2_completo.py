#!/usr/bin/env python3
"""
PPTX 2 COMPLETO: Lección extensa de Gestión y Gobernanza de Datos
con la plataforma Data Lake Governance Hub como caso práctico.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== COLORES =====
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
SLATE_900 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_800 = RGBColor(0x1E, 0x40, 0x5F)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x05, 0x96, 0x69)
PINK = RGBColor(0xEC, 0x48, 0x99)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

def set_slide_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return tb

def add_lines(slide, left, top, width, height, lines, size=14, color=SLATE_300, spacing=4):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1] if len(line) > 1 else color
            p.font.bold = line[2] if len(line) > 2 else False
            p.font.size = Pt(line[3] if len(line) > 3 else size)
        else:
            p.text = line
            p.font.color.rgb = color
            p.font.size = Pt(size)
        p.font.name = 'Calibri'
        p.space_after = Pt(spacing)
    return tb

def add_bullets(slide, left, top, width, height, items, size=14, color=SLATE_300):
    return add_lines(slide, left, top, width, height, items, size, color, spacing=5)

def add_code(slide, left, top, width, height, code, size=10):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p.font.name = 'Consolas'
    return sh

def badge(slide, left, top, text, bg, fg=WHITE, w=1.5, h=0.3):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'

def title_slide(prs, title, subtitle, extra=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_text(s, 0.5, 1.5, 9, 2, title, 36, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, 0.5, 3.8, 9, 1, subtitle, 16, SLATE_400, False, PP_ALIGN.CENTER)
    if extra:
        add_text(s, 0.5, 5.0, 9, 0.5, extra, 13, BLUE, False, PP_ALIGN.CENTER)
    add_text(s, 0.5, 6.8, 9, 0.4, "Gestión y Gobernanza de Datos — Universidad EAFIT — 2026", 10, SLATE_500, False, PP_ALIGN.CENTER)
    return s

def section(prs, num, title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, SLATE_900)
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.3), Inches(1.6), Inches(0.8), Inches(0.8))
    sh.fill.solid()
    sh.fill.fore_color.rgb = BLUE
    sh.line.fill.background()
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_text(s, 0.5, 2.7, 9, 1, title, 30, WHITE, True, PP_ALIGN.CENTER)
    if sub:
        add_text(s, 0.5, 3.7, 9, 0.8, sub, 14, SLATE_400, False, PP_ALIGN.CENTER)
    return s

def slide(prs, title, bdg=None, bdg_color=BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_text(s, 0.5, 0.25, 7.5, 0.6, title, 21, WHITE, True)
    if bdg:
        badge(s, 8.2, 0.28, bdg, bdg_color, w=1.3)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.82), Inches(9), Inches(0.015))
    sh.fill.solid()
    sh.fill.fore_color.rgb = SLATE_800
    sh.line.fill.background()
    return s

# ================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # PORTADA
    # ================================================================
    title_slide(prs,
        "Gestión y Gobernanza de Datos\nen la Práctica",
        "Lección completa: del marco teórico DAMA-DMBOK2\nal laboratorio real con Data Lake Governance Hub",
        "datalake-governance-hub.vercel.app")

    # ================================================================
    # AGENDA
    # ================================================================
    s = slide(prs, "Contenido de la lección")
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("Parte I — Fundamentos de Gestión y Gobernanza de Datos", BLUE, True, 16),
        ("   1. ¿Qué es la Gestión de Datos? El framework DAMA-DMBOK2", SLATE_300, False, 12),
        ("   2. ¿Qué es la Gobernanza de Datos? Roles, políticas y procesos", SLATE_300, False, 12),
        ("   3. Marco regulatorio: Habeas Data, GDPR, y por qué importa", SLATE_300, False, 12),
        ("   4. Data Lake: arquitectura, buenas prácticas y anti-patrones", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Parte II — La plataforma como caso de estudio", GREEN, True, 16),
        ("   5. Capa Bronze: ingestión y datos crudos", SLATE_300, False, 12),
        ("   6. Capa Silver: transformación, calidad y validación", SLATE_300, False, 12),
        ("   7. Capa Gold: datos de negocio y consumo", SLATE_300, False, 12),
        ("   8. Catálogo de datos: metadatos y descubrimiento", SLATE_300, False, 12),
        ("   9. Calidad de datos: las 6 dimensiones del DMBOK2", SLATE_300, False, 12),
        ("   10. Linaje: trazabilidad de origen a consumo", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Parte III — Seguridad y operación", RED, True, 16),
        ("   11. Seguridad: RLS, enmascaramiento PII, RBAC", SLATE_300, False, 12),
        ("   12. Auditoría: cumplimiento y trazabilidad de accesos", SLATE_300, False, 12),
        ("   13. Operación: guía tab por tab de la plataforma", SLATE_300, False, 12),
        ("   14. Modelo de madurez y siguientes pasos", SLATE_300, False, 12),
    ], spacing=2)

    # ================================================================
    # PARTE I: FUNDAMENTOS
    # ================================================================

    # --- MÓDULO 1: GESTIÓN DE DATOS ---
    section(prs, 1, "¿Qué es la Gestión de Datos?", "El framework DAMA-DMBOK2 y sus 11 áreas de conocimiento")

    s = slide(prs, "Gestión de Datos: definición", "DAMA-DMBOK2", BLUE)
    add_text(s, 0.5, 1.0, 9, 0.8, "La Gestión de Datos (Data Management) es el desarrollo, ejecución y supervisión de planes, políticas, programas y prácticas que entregan, controlan, protegen y mejoran el valor de los datos e información.", 14, SLATE_300)
    add_text(s, 0.5, 2.0, 9, 0.3, "— DAMA International, DMBOK2 (2017)", 12, SLATE_500)
    add_text(s, 0.5, 2.7, 9, 0.4, "¿Por qué importa?", 15, AMBER, True)
    add_bullets(s, 0.5, 3.1, 9, 3, [
        "Los datos son el activo más valioso de una organización moderna — pero solo si se gestionan bien",
        "Sin gestión: datos duplicados, inconsistentes, inaccesibles, inseguros → decisiones incorrectas",
        "Con gestión: datos confiables, accesibles, seguros, trazables → ventaja competitiva",
        "IDC estima que las empresas pierden hasta el 30% de ingresos por mala calidad de datos",
        "El mercado de Data Governance pasará de $2.1B (2023) a $7.9B (2028) — CAGR del 25%",
    ], 13)
    add_text(s, 0.5, 6.0, 9, 0.8, "En nuestra plataforma: implementamos 6 de las 11 áreas de conocimiento del DMBOK2 de forma real y operativa.", 12, GREEN)

    # DAMA Wheel
    s = slide(prs, "Las 11 áreas de conocimiento del DAMA-DMBOK2", "DMBOK2", BLUE)
    add_text(s, 0.5, 1.0, 9, 0.4, "El DAMA-DMBOK2 organiza la gestión de datos en 11 áreas. La Gobernanza está en el centro:", 13, SLATE_300)

    areas = [
        ("1. Gobernanza de Datos", "Centro del framework — dirige todas las demás", RED, True),
        ("2. Arquitectura de Datos", "Modelos, estándares, Blueprint → Medallion en la plataforma", BLUE, False),
        ("3. Modelado y Diseño", "Modelos conceptuales, lógicos, físicos → esquemas bronze/silver/gold", BLUE, False),
        ("4. Almacenamiento y Operaciones", "Bases de datos, backups, rendimiento → Supabase PostgreSQL", CYAN, False),
        ("5. Seguridad de Datos", "Acceso, cifrado, privacidad → RLS, enmascaramiento, RBAC", RED, False),
        ("6. Integración e Interoperabilidad", "ETL, ELT, APIs, migración → Pipeline Bronze→Silver→Gold", GREEN, False),
        ("7. Gestión Documental y de Contenido", "Documentos, archivos, contenido no estructurado", SLATE_400, False),
        ("8. Datos Maestros y de Referencia (MDM)", "Golden record, entidades únicas, dominios", SLATE_400, False),
        ("9. Data Warehousing y BI", "Data Warehouse, Data Lake, analytics → Dashboard Gold", GREEN, False),
        ("10. Metadatos", "Catálogo, linaje, glosario → Tab Catálogo y Tab Linaje", GREEN, False),
        ("11. Calidad de Datos", "Dimensiones, reglas, monitoreo → Tab Calidad, quality_score", GREEN, False),
    ]
    add_lines(s, 0.5, 1.5, 9, 5.5, [
        (f"{a[0]}", a[2], a[3], 12) if a[3] else (f"  {a[0]}: {a[1]}", a[2], False, 11)
        for a in areas
    ], spacing=2)
    add_text(s, 0.5, 6.6, 9, 0.4, "Verde = implementado en la plataforma  |  Gris = no aplica directamente al ejercicio", 10, SLATE_500)

    # Áreas implementadas en detalle
    s = slide(prs, "¿Cómo implementamos cada área DAMA en la plataforma?", "PRÁCTICA", GREEN)
    add_lines(s, 0.5, 1.0, 4.3, 6.2, [
        ("Gobernanza (área 1)", RED, True, 13),
        ("  Catálogo obligatorio, clasificación,", SLATE_300, False, 11),
        ("  owners por dataset, políticas de acceso", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Arquitectura (área 2)", BLUE, True, 13),
        ("  Medallion: Bronze → Silver → Gold", SLATE_300, False, 11),
        ("  Separación en esquemas PostgreSQL", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Seguridad (área 5)", RED, True, 13),
        ("  RLS por rol, PII masking, RBAC,", SLATE_300, False, 11),
        ("  auditoría de accesos, cifrado", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Integración (área 6)", GREEN, True, 13),
        ("  Pipeline SQL: JSONB → tipado → agregado", SLATE_300, False, 11),
        ("  API REST Supabase → Next.js frontend", SLATE_300, False, 11),
    ], spacing=2)
    add_lines(s, 5, 1.0, 4.5, 6.2, [
        ("BI / Data Lake (área 9)", CYAN, True, 13),
        ("  Dashboard Gold con KPIs de negocio", SLATE_300, False, 11),
        ("  Vistas materializadas para rendimiento", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Metadatos (área 10)", PURPLE, True, 13),
        ("  gold.data_catalog: dueño, clasificación,", SLATE_300, False, 11),
        ("  PII, freshness, quality_score por dataset", SLATE_300, False, 11),
        ("  Linaje: source_record_id en Silver", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Calidad (área 11)", AMBER, True, 13),
        ("  6 dimensiones DAMA implementadas", SLATE_300, False, 11),
        ("  7 reglas de calidad automatizadas", SLATE_300, False, 11),
        ("  Quality score por registro y dataset", SLATE_300, False, 11),
        ("  Dashboard de calidad en tiempo real", SLATE_300, False, 11),
    ], spacing=2)

    # --- MÓDULO 2: GOBERNANZA ---
    section(prs, 2, "¿Qué es la Gobernanza de Datos?", "Roles, políticas, procesos y la diferencia con gestión")

    s = slide(prs, "Gobernanza vs Gestión de Datos", "DIFERENCIA", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.5, "No son lo mismo. La Gobernanza DIRIGE, la Gestión EJECUTA.", 15, WHITE, True)
    add_lines(s, 0.5, 1.7, 4.3, 3, [
        ("Gobernanza de Datos", PURPLE, True, 15),
        ("  Qué: Define políticas, estándares y roles", SLATE_300, False, 12),
        ("  Quién: Comité de gobernanza, Data Owners", SLATE_300, False, 12),
        ("  Cómo: Políticas, frameworks, auditorías", SLATE_300, False, 12),
        ("  Cuándo: Estratégico, a largo plazo", SLATE_300, False, 12),
        ("  Ejemplo: 'Los datos PII deben estar", SLATE_300, False, 12),
        ("  enmascarados para roles no-admin'", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 5, 1.7, 4.5, 3, [
        ("Gestión de Datos", BLUE, True, 15),
        ("  Qué: Implementa las políticas definidas", SLATE_300, False, 12),
        ("  Quién: Data Engineers, DBAs, Stewards", SLATE_300, False, 12),
        ("  Cómo: Herramientas, pipelines, código", SLATE_300, False, 12),
        ("  Cuándo: Operativo, día a día", SLATE_300, False, 12),
        ("  Ejemplo: 'Crear vista ventas_masked", SLATE_300, False, 12),
        ("  con CONCAT(LEFT(nit,3),****)'", SLATE_300, False, 12),
    ], spacing=3)
    add_text(s, 0.5, 5.2, 9, 0.4, "En nuestra plataforma:", 14, GREEN, True)
    add_bullets(s, 0.5, 5.6, 9, 1.5, [
        "La Gobernanza dice: 'Todo dataset debe tener dueño, clasificación y quality score' → gold.data_catalog",
        "La Gestión implementa: el pipeline SQL que calcula quality_score y lo registra en el catálogo",
    ], 12)

    # Roles
    s = slide(prs, "Roles en la Gobernanza de Datos", "ROLES", CYAN)
    add_text(s, 0.5, 1.0, 9, 0.4, "Cada rol tiene responsabilidades claras. En la plataforma, se mapean a los 4 usuarios:", 13, SLATE_300)

    add_lines(s, 0.5, 1.6, 4.3, 5.5, [
        ("Data Owner (Dueño de datos)", WHITE, True, 14),
        ("  Responsable de un dominio de datos", SLATE_300, False, 11),
        ("  Aprueba quién accede a sus datos", SLATE_300, False, 11),
        ("  Define reglas de calidad y retención", SLATE_300, False, 11),
        ("  En plataforma: 'equipo-datos' como owner", SLATE_300, False, 11),
        ("  de ventas_raw y ventas_clean", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Data Steward (Custodio de datos)", WHITE, True, 14),
        ("  Operacionaliza las políticas del Owner", SLATE_300, False, 11),
        ("  Monitorea calidad, resuelve incidentes", SLATE_300, False, 11),
        ("  Mantiene el catálogo actualizado", SLATE_300, False, 11),
        ("  En plataforma: quien actualiza data_catalog", SLATE_300, False, 11),
    ], spacing=2)

    add_lines(s, 5, 1.6, 4.5, 5.5, [
        ("Data Engineer (Ingeniero de datos)", WHITE, True, 14),
        ("  Construye y mantiene los pipelines", SLATE_300, False, 11),
        ("  Implementa las transformaciones B→S→G", SLATE_300, False, 11),
        ("  En plataforma: Carlos Martínez", SLATE_300, False, 11),
        ("  Acceso: Bronze R/W + Silver R/W + Gold R", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Data Consumer (Consumidor)", WHITE, True, 14),
        ("  Usa los datos para análisis o decisiones", SLATE_300, False, 11),
        ("  Data Scientist: María López", SLATE_300, False, 11),
        ("    → Silver masked + Gold + Sandbox", SLATE_300, False, 11),
        ("  Analyst: Pedro García", SLATE_300, False, 11),
        ("    → Solo Gold (lectura)", SLATE_300, False, 11),
    ], spacing=2)

    # Políticas
    s = slide(prs, "Políticas de Gobernanza implementadas", "POLÍTICAS", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.4, "Una política es una regla organizacional. Sin políticas, la gobernanza es solo palabras.", 13, SLATE_300)

    add_lines(s, 0.5, 1.6, 9, 5.5, [
        ("Política 1: Clasificación obligatoria de datos", WHITE, True, 14),
        ("  Todo dataset debe clasificarse como public, internal, confidential o restricted.", SLATE_300, False, 12),
        ("  En plataforma: columna classification en gold.data_catalog. Sin clasificación, no entra al Lake.", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("Política 2: PII identificado y protegido", WHITE, True, 14),
        ("  Datos con información personal (NIT, nombre, email) deben marcarse como has_pii=true.", SLATE_300, False, 12),
        ("  En plataforma: bronze.ventas_raw tiene PII=Sí. Se enmascara en silver.ventas_masked.", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("Política 3: Mínimo privilegio (Least Privilege)", WHITE, True, 14),
        ("  Cada usuario accede SOLO a lo que necesita. Un analista no necesita ver Bronze.", SLATE_300, False, 12),
        ("  En plataforma: RLS en Silver, analista solo ve Gold, data scientist ve Silver masked.", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("Política 4: Auditoría completa de accesos", WHITE, True, 14),
        ("  Todo acceso al Lake queda registrado: quién, qué, cuándo, cuánto, desde dónde.", SLATE_300, False, 12),
        ("  En plataforma: audit.access_log con timestamp, email, rol, acción, tabla, filas, IP.", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("Política 5: Calidad medida y monitoreada", WHITE, True, 14),
        ("  Cada dataset tiene un quality_score. Si baja de 80%, se genera alerta.", SLATE_300, False, 12),
        ("  En plataforma: quality_score por registro en Silver + por dataset en data_catalog.", SLATE_400, False, 11),
    ], spacing=1)

    # --- MÓDULO 3: REGULACIÓN ---
    section(prs, 3, "Marco Regulatorio", "¿Por qué la gobernanza de datos NO es opcional?")

    s = slide(prs, "Regulaciones que exigen gobernanza de datos", "LEGAL", RED)
    add_lines(s, 0.5, 1.0, 4.3, 6, [
        ("Colombia", AMBER, True, 15),
        ("", WHITE, False, 2),
        ("Ley 1581 de 2012 (Habeas Data)", WHITE, True, 13),
        ("  Protección de datos personales.", SLATE_300, False, 11),
        ("  Requiere: consentimiento, finalidad,", SLATE_300, False, 11),
        ("  acceso, rectificación, eliminación.", SLATE_300, False, 11),
        ("  En plataforma: PII masking, clasificación", SLATE_400, False, 10),
        ("", WHITE, False, 4),
        ("Ley 1266 de 2008", WHITE, True, 13),
        ("  Datos financieros y crediticios.", SLATE_300, False, 11),
        ("  En plataforma: classification='confidential'", SLATE_400, False, 10),
        ("", WHITE, False, 4),
        ("Circular 007/2018 SFC", WHITE, True, 13),
        ("  Ciberseguridad en sector financiero.", SLATE_300, False, 11),
        ("  En plataforma: auditoría, cifrado, RLS", SLATE_400, False, 10),
    ], spacing=2)

    add_lines(s, 5, 1.0, 4.5, 6, [
        ("Internacional", BLUE, True, 15),
        ("", WHITE, False, 2),
        ("GDPR (Unión Europea, 2018)", WHITE, True, 13),
        ("  El estándar global de privacidad.", SLATE_300, False, 11),
        ("  Multas: hasta 4% de ingresos globales.", SLATE_300, False, 11),
        ("  Requiere: DPO, DPIA, portabilidad,", SLATE_300, False, 11),
        ("  derecho al olvido, breach notification.", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("SOX (Sarbanes-Oxley, EE.UU.)", WHITE, True, 13),
        ("  Control de datos financieros.", SLATE_300, False, 11),
        ("  Requiere: auditoría, integridad, acceso.", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("HIPAA (EE.UU. — Salud)", WHITE, True, 13),
        ("  Datos de salud protegidos.", SLATE_300, False, 11),
        ("  Requiere: cifrado, acceso controlado.", SLATE_300, False, 11),
    ], spacing=2)

    add_text(s, 0.5, 6.8, 9, 0.4, "La plataforma implementa controles que satisfacen requisitos de Habeas Data, GDPR y SOX.", 11, GREEN)

    # Casos de fallo
    s = slide(prs, "Cuando la gobernanza falla: casos reales", "CASOS", RED)
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("Capital One (2019) — 100 millones de registros filtrados", RED, True, 15),
        ("  Un bucket S3 mal configurado (sin control de acceso) permitió que un atacante", SLATE_300, False, 12),
        ("  accediera a nombres, NITs, scores crediticios. Multa: $80 millones USD.", SLATE_300, False, 12),
        ("  Lección: NUNCA dejar un Data Lake sin control de acceso. → En nuestra plataforma: RLS + RBAC.", AMBER, False, 11),
        ("", WHITE, False, 6),
        ("Equifax (2017) — 147 millones de personas afectadas", RED, True, 15),
        ("  Datos de crédito expuestos por una vulnerabilidad no parcheada + falta de monitoreo.", SLATE_300, False, 12),
        ("  Multa: $700 millones USD. El CEO fue despedido.", SLATE_300, False, 12),
        ("  Lección: Monitorear, auditar, cifrar. → En nuestra plataforma: auditoría + cifrado Supabase.", AMBER, False, 11),
        ("", WHITE, False, 6),
        ("Uber (2016) — 57 millones de registros", RED, True, 15),
        ("  Un ingeniero subió credenciales de la BD a un repo público. Uber pagó $100K a los atacantes", SLATE_300, False, 12),
        ("  para que eliminaran los datos en vez de reportar. Multa: $148 millones USD.", SLATE_300, False, 12),
        ("  Lección: Variables sensibles NUNCA en el código. → .env.local en .gitignore.", AMBER, False, 11),
        ("", WHITE, False, 6),
        ("Facebook/Cambridge Analytica (2018) — 87 millones de perfiles", RED, True, 15),
        ("  Datos compartidos sin consentimiento para microtargeting político.", SLATE_300, False, 12),
        ("  Lección: Clasificación de datos + consentimiento + auditoría de quién accede a qué.", AMBER, False, 11),
    ], spacing=1)

    # --- MÓDULO 4: DATA LAKE ---
    section(prs, 4, "Data Lake: Arquitectura y Mejores Prácticas", "Qué es, cómo se organiza, y qué evitar")

    s = slide(prs, "Data Lake vs Data Warehouse vs Data Lakehouse", "COMPARACIÓN", BLUE)
    add_lines(s, 0.5, 1.0, 3, 5.5, [
        ("Data Warehouse", AMBER, True, 14),
        ("", WHITE, False, 2),
        ("Solo datos estructurados", SLATE_300, False, 11),
        ("Schema-on-write", SLATE_300, False, 11),
        ("SQL rápido, optimizado", SLATE_300, False, 11),
        ("Costoso por GB", SLATE_300, False, 11),
        ("Usuarios: analistas BI", SLATE_300, False, 11),
        ("Ejemplos: Snowflake,", SLATE_300, False, 11),
        ("  Redshift, BigQuery", SLATE_300, False, 11),
    ], spacing=3)
    add_lines(s, 3.5, 1.0, 3, 5.5, [
        ("Data Lake", BLUE, True, 14),
        ("", WHITE, False, 2),
        ("Todos los formatos", SLATE_300, False, 11),
        ("Schema-on-read", SLATE_300, False, 11),
        ("Flexible, escalable", SLATE_300, False, 11),
        ("Barato ($0.02/GB/mes)", SLATE_300, False, 11),
        ("Usuarios: data eng/sci", SLATE_300, False, 11),
        ("Ejemplos: S3+Glue,", SLATE_300, False, 11),
        ("  ADLS, GCS", SLATE_300, False, 11),
    ], spacing=3)
    add_lines(s, 6.5, 1.0, 3.3, 5.5, [
        ("Data Lakehouse", GREEN, True, 14),
        ("", WHITE, False, 2),
        ("Lo mejor de ambos", SLATE_300, False, 11),
        ("Schema enforcement", SLATE_300, False, 11),
        ("ACID + Time Travel", SLATE_300, False, 11),
        ("Barato + rápido", SLATE_300, False, 11),
        ("Usuarios: todos", SLATE_300, False, 11),
        ("Ejemplos: Databricks,", SLATE_300, False, 11),
        ("  Delta Lake, Iceberg", SLATE_300, False, 11),
    ], spacing=3)
    add_text(s, 0.5, 5.8, 9, 0.8, "Nuestra plataforma implementa un Data Lake con arquitectura Medallion en PostgreSQL. En producción, se usaría S3/ADLS/GCS + Spark/dbt + Athena/BigQuery.", 12, SLATE_400)

    # Medallion
    s = slide(prs, "Arquitectura Medallion: Bronze → Silver → Gold", "MEDALLION", BLUE)
    add_text(s, 0.5, 1.0, 9, 0.4, "El estándar de la industria para organizar un Data Lake. Popularizado por Databricks (2020).", 13, SLATE_300)

    badge(s, 0.5, 1.6, "BRONZE", AMBER, w=2.8)
    add_lines(s, 0.5, 2.0, 2.8, 2.2, [
        ("Zona de aterrizaje", WHITE, True, 12),
        ("Datos TAL CUAL llegan", SLATE_300, False, 11),
        ("JSONB, CSV, Parquet", SLATE_300, False, 11),
        ("Inmutables (append-only)", SLATE_300, False, 11),
        ("Con errores, duplicados", SLATE_300, False, 11),
        ("Retención: 1-7 años", SLATE_300, False, 11),
    ], spacing=2)
    add_text(s, 3.35, 2.5, 0.5, 0.5, "→", 28, SLATE_400, align=PP_ALIGN.CENTER)

    badge(s, 3.8, 1.6, "SILVER", PURPLE, w=2.5)
    add_lines(s, 3.8, 2.0, 2.5, 2.2, [
        ("Zona conformada", WHITE, True, 12),
        ("Tipado, validado", SLATE_300, False, 11),
        ("Deduplicado, normalizado", SLATE_300, False, 11),
        ("Quality score calculado", SLATE_300, False, 11),
        ("Formato columnar (Parquet)", SLATE_300, False, 11),
        ("Linaje a Bronze", SLATE_300, False, 11),
    ], spacing=2)
    add_text(s, 6.35, 2.5, 0.5, 0.5, "→", 28, SLATE_400, align=PP_ALIGN.CENTER)

    badge(s, 6.8, 1.6, "GOLD", EMERALD, w=2.7)
    add_lines(s, 6.8, 2.0, 2.7, 2.2, [
        ("Zona de negocio", WHITE, True, 12),
        ("Agregaciones, KPIs", SLATE_300, False, 11),
        ("Vistas materializadas", SLATE_300, False, 11),
        ("Modelo dimensional", SLATE_300, False, 11),
        ("Acceso controlado por rol", SLATE_300, False, 11),
        ("SLAs de calidad", SLATE_300, False, 11),
    ], spacing=2)

    add_text(s, 0.5, 4.5, 9, 0.4, "En nuestra plataforma (datos reales):", 14, GREEN, True)
    add_code(s, 0.5, 4.9, 9, 2.2,
"""Bronze: bronze.ventas_raw      → 8 registros JSONB (con errores intencionales)
        bronze.sensores_raw    → 5 lecturas IoT (con anomalía -5.0 L/s)

Silver: silver.ventas_clean     → 7 registros tipados (1 rechazado por NULL)
        silver.sensores_clean  → 5 lecturas validadas (anomalía marcada)
        silver.ventas_masked   → Vista con PII enmascarado

Gold:   gold.ventas_por_ciudad  → Agregación por ciudad (vista materializada)
        gold.ventas_por_producto→ Agregación por producto
        gold.estado_sensores   → Dashboard IoT con anomalías
        gold.data_catalog      → Catálogo de metadatos""", 10)

    # Anti-patrones
    s = slide(prs, "Anti-patrones: cómo NO gestionar un Data Lake", "ANTI-PATRONES", RED)
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("1. Data Swamp — tirar todo sin estructura", RED, True, 14),
        ("   archivos como 'reporte_final_v2_FIXED_backup.csv' sin metadatos ni organización.", SLATE_300, False, 12),
        ("   → Solución: arquitectura Medallion + catálogo obligatorio (como nuestra plataforma).", GREEN, False, 11),
        ("", WHITE, False, 4),
        ("2. Todos son admin — sin control de acceso", RED, True, 14),
        ("   Un becario puede borrar Bronze. Un analista ve NITs reales sin necesitarlos.", SLATE_300, False, 12),
        ("   → Solución: RBAC con 4 roles + RLS + enmascaramiento PII.", GREEN, False, 11),
        ("", WHITE, False, 4),
        ("3. JSON everywhere — no convertir a formato columnar", RED, True, 14),
        ("   Consultas 10x más lentas, almacenamiento 5x más caro. Athena cobra por byte escaneado.", SLATE_300, False, 12),
        ("   → Solución: Parquet desde Silver. En nuestra plataforma: columnas tipadas con CHECK.", GREEN, False, 11),
        ("", WHITE, False, 4),
        ("4. Set it and forget it — sin monitoreo", RED, True, 14),
        ("   El pipeline se rompe en febrero, nadie se da cuenta hasta abril. Dashboard con datos viejos.", SLATE_300, False, 12),
        ("   → Solución: quality_score + freshness en catálogo + alertas cuando calidad baja de 80%.", GREEN, False, 11),
        ("", WHITE, False, 4),
        ("5. Sin linaje — no saber de dónde viene el dato", RED, True, 14),
        ("   Un KPI parece incorrecto, pero nadie puede rastrear el origen para verificar.", SLATE_300, False, 12),
        ("   → Solución: source_record_id en Silver + tab Linaje con trazabilidad registro por registro.", GREEN, False, 11),
    ], spacing=1)

    # ================================================================
    # PARTE II: LA PLATAFORMA
    # ================================================================

    # --- MÓDULO 5: BRONZE ---
    section(prs, 5, "Capa Bronze en la plataforma", "Tab 'Bronze (Raw)' — los datos crudos e inmutables")

    s = slide(prs, "Bronze: qué verás en la plataforma", "BRONZE", AMBER)
    add_text(s, 0.5, 1.0, 9, 0.5, "Al abrir la tab 'Bronze (Raw)' ves 8 tarjetas, cada una es un registro JSONB tal cual llegó del ERP.", 14, SLATE_300)
    add_text(s, 0.5, 1.7, 9, 0.4, "Cada tarjeta muestra:", 14, AMBER, True)
    add_bullets(s, 0.5, 2.1, 9, 1.5, [
        "ID del registro, sistema fuente (SAP-ERP), archivo original (export_ventas_20260301.json)",
        "El JSON completo expandido con todos los campos",
        "Etiquetas rojas/amarillas/azules cuando se detectan problemas de calidad",
    ], 13)
    add_text(s, 0.5, 3.5, 9, 0.4, "Los 3 errores intencionales (buscarlos es parte del ejercicio):", 14, RED, True)
    add_code(s, 0.5, 3.9, 4.3, 1,
"""Registro #4:
  "cantidad": null    ← NULL
  "vendedor": ""      ← vacío
  "ciudad": "medellin"← sin normalizar""", 11)
    add_lines(s, 5, 3.9, 4.5, 1.5, [
        ("¿Por qué errores intencionales?", AMBER, True, 13),
        ("  En la vida real, los sistemas fuente", SLATE_300, False, 11),
        ("  SIEMPRE envían datos sucios.", SLATE_300, False, 11),
        ("  El Lake debe saber manejarlos.", SLATE_300, False, 11),
    ], spacing=3)
    add_text(s, 0.5, 5.3, 9, 0.4, "Regla de oro de Bronze:", 14, GREEN, True)
    add_text(s, 0.5, 5.7, 9, 0.8, "NUNCA modificar Bronze. Si hay errores, se corrigen en Silver. Bronze es evidencia del dato original — si Silver tiene un bug, puedes reprocesar desde Bronze. Si borras Bronze, pierdes la fuente de verdad.", 13, SLATE_300)

    # Bronze SQL
    s = slide(prs, "Bronze: la estructura SQL detrás", "SQL", AMBER)
    add_text(s, 0.5, 1.0, 9, 0.4, "¿Por qué JSONB y no columnas tipadas en Bronze?", 14, AMBER, True)
    add_bullets(s, 0.5, 1.4, 9, 1.5, [
        "JSONB acepta CUALQUIER estructura — si el ERP agrega un campo 'descuento', no rompe la tabla",
        "Es schema-on-read: guardas primero, defines esquema después (en Silver)",
        "Puedes almacenar datos de fuentes completamente diferentes en la misma tabla",
    ], 12)
    add_code(s, 0.5, 3.0, 9, 2.5,
"""-- Así se carga un dato a Bronze (la plataforma lo hizo por ti):
INSERT INTO bronze.ventas_raw (raw_data, source_system, ingested_by, file_name)
VALUES (
  '{"fecha":"2026-03-01","producto":"Sensor IoT X200",
    "cantidad":150,"precio_unit":45000,
    "vendedor":"Juan Pérez","ciudad":"Medellín",
    "cliente_nit":"900123456","cliente_nombre":"Acueducto Municipal"}',
  'SAP-ERP',           -- ¿De qué sistema viene?
  'etl-pipeline',      -- ¿Quién lo cargó?
  'export_ventas_20260301.json'  -- ¿De qué archivo?
);""", 11)
    add_text(s, 0.5, 5.7, 9, 0.8, "Metadatos de ingestión: source_system, ingested_by, file_name y batch_id permiten rastrear cada carga. Si un lote tiene problemas, puedes identificar y reprocesar solo ese batch.", 12, SLATE_400)

    # --- MÓDULO 6: SILVER ---
    section(prs, 6, "Capa Silver: transformación y calidad", "Tabs 'Silver (Clean)' y 'Bronze vs Silver'")

    s = slide(prs, "Las 6 transformaciones Bronze → Silver", "ETL", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.4, "Cada transformación resuelve un problema específico de calidad:", 13, SLATE_300)

    add_text(s, 0.5, 1.6, 4.3, 0.3, "1. Casteo de tipos", 13, WHITE, True)
    add_code(s, 0.5, 1.9, 4.3, 0.5, "(raw_data->>'fecha')::date\n(raw_data->>'cantidad')::integer", 10)
    add_text(s, 0.5, 2.4, 4.3, 0.3, "Texto → tipo correcto. Permite filtros y math.", 10, SLATE_500)

    add_text(s, 5, 1.6, 4.5, 0.3, "2. Tratamiento de vacíos", 13, WHITE, True)
    add_code(s, 5, 1.9, 4.5, 0.5, "COALESCE(NULLIF(vendedor,''),\n         'Sin asignar')", 10)
    add_text(s, 5, 2.4, 4.5, 0.3, "String vacío → valor por defecto.", 10, SLATE_500)

    add_text(s, 0.5, 2.9, 4.3, 0.3, "3. Normalización", 13, WHITE, True)
    add_code(s, 0.5, 3.2, 4.3, 0.5, "INITCAP(raw_data->>'ciudad')\n-- 'medellin' → 'Medellin'", 10)
    add_text(s, 0.5, 3.7, 4.3, 0.3, "Consistencia para GROUP BY.", 10, SLATE_500)

    add_text(s, 5, 2.9, 4.5, 0.3, "4. Filtrado de inválidos", 13, WHITE, True)
    add_code(s, 5, 3.2, 4.5, 0.5, "WHERE cantidad IS NOT NULL\n  AND cantidad > 0", 10)
    add_text(s, 5, 3.7, 4.5, 0.3, "Rechaza registros imposibles.", 10, SLATE_500)

    add_text(s, 0.5, 4.2, 4.3, 0.3, "5. Quality Score", 13, WHITE, True)
    add_code(s, 0.5, 4.5, 4.3, 0.5, "CASE WHEN vendedor=''\n  THEN 0.70 ELSE 1.00 END", 10)
    add_text(s, 0.5, 5.0, 4.3, 0.3, "Mide confiabilidad por registro.", 10, SLATE_500)

    add_text(s, 5, 4.2, 4.5, 0.3, "6. Campo calculado", 13, WHITE, True)
    add_code(s, 5, 4.5, 4.5, 0.5, "ingreso_total GENERATED ALWAYS\n  AS (cantidad*precio_unit) STORED", 10)
    add_text(s, 5, 5.0, 4.5, 0.3, "Calculado automáticamente, sin errores.", 10, SLATE_500)

    add_text(s, 0.5, 5.5, 9, 0.4, "Resultado: 8 Bronze → 7 Silver (1 rechazado: registro #4 por cantidad NULL)", 14, GREEN, True)
    add_text(s, 0.5, 6.0, 9, 0.8, "El registro rechazado permanece en Bronze como evidencia. En la tab 'Bronze vs Silver' puedes verlo marcado con etiqueta roja 'RECHAZADO' y el motivo específico.", 12, SLATE_400)

    # Bronze vs Silver
    s = slide(prs, "Tab 'Bronze vs Silver': interpretar la comparación", "COMPARAR", CYAN)
    add_text(s, 0.5, 1.0, 9, 0.4, "Esta tab muestra cada registro Bronze lado a lado con su versión Silver.", 14, SLATE_300)
    add_text(s, 0.5, 1.5, 9, 0.4, "Cómo leer cada tarjeta:", 14, CYAN, True)
    add_bullets(s, 0.5, 1.9, 9, 2.5, [
        "Encabezado: B#1 → S#1 (Bronze ID 1 se transformó en Silver ID 1)",
        "Etiqueta VERDE 'LIMPIO': pasó sin problemas, quality_score = 1.00",
        "Etiqueta AMARILLA 'N transformaciones': se corrigieron N campos (vendedor vacío, ciudad)",
        "Etiqueta ROJA 'RECHAZADO': NO pasó a Silver — con motivo (cantidad NULL)",
        "Izquierda: JSON de Bronze con campos problemáticos resaltados en ROJO",
        "Derecha: columnas de Silver con campos corregidos resaltados en VERDE",
        "Footer: detalle de cada transformación aplicada",
    ], 12)
    add_text(s, 0.5, 5.0, 9, 0.4, "Ejercicio:", 14, AMBER, True)
    add_bullets(s, 0.5, 5.4, 9, 1.5, [
        "Abre 'Bronze vs Silver'. Encuentra el registro RECHAZADO. ¿Cuántos problemas tenía? ¿Cuáles?",
        "Busca un registro con etiqueta amarilla. ¿Qué campo se transformó? ¿Cómo cambió?",
    ], 12)

    # --- MÓDULO 7: GOLD ---
    section(prs, 7, "Capa Gold: datos listos para negocio", "Tab 'Dashboard Gold'")

    s = slide(prs, "Gold: vistas materializadas y KPIs", "GOLD", GREEN)
    add_text(s, 0.5, 1.0, 9, 0.5, "Gold no almacena datos nuevos. Resume y agrega Silver para que el negocio tome decisiones.", 14, SLATE_300)
    add_code(s, 0.5, 1.6, 9, 2.3,
"""-- Vista materializada = consulta SQL pre-calculada guardada como tabla
CREATE MATERIALIZED VIEW gold.ventas_por_ciudad AS
SELECT
    ciudad,
    COUNT(*)           AS num_transacciones,    -- ¿Cuántas ventas?
    SUM(cantidad)      AS unidades_vendidas,    -- ¿Cuántas unidades?
    SUM(ingreso_total) AS ingreso_total,        -- ¿Cuánto dinero?
    ROUND(AVG(ingreso_total), 2) AS ticket_promedio,
    ROUND(AVG(quality_score), 2) AS calidad_promedio  -- ¿Qué tan confiable?
FROM silver.ventas_clean    -- SOLO lee Silver, NUNCA Bronze
GROUP BY ciudad ORDER BY ingreso_total DESC;

-- Para actualizar después de nuevos datos en Silver:
REFRESH MATERIALIZED VIEW gold.ventas_por_ciudad;""", 10)
    add_text(s, 0.5, 4.2, 9, 0.4, "¿Por qué el dashboard SOLO lee Gold?", 14, AMBER, True)
    add_bullets(s, 0.5, 4.6, 9, 2.5, [
        "Rendimiento: Gold ya tiene las agregaciones calculadas → consultas instantáneas",
        "Seguridad: un analista no necesita (ni debe) ver datos individuales con PII",
        "Separación: si Silver se reprocesa, Gold no se afecta hasta que se refresque",
        "Gobernanza: el código dice supabase.from('gold_ventas_por_ciudad') — auditable",
    ], 12)

    # Dashboard detalles
    s = slide(prs, "Tab 'Dashboard Gold': los 4 KPIs explicados", "DASHBOARD", GREEN)
    add_lines(s, 0.5, 1.0, 4.3, 3, [
        ("KPI 1: Ingreso Total", GREEN, True, 14),
        ("  SUM(cantidad × precio) de todas las ventas", SLATE_300, False, 12),
        ("  que pasaron de Bronze a Silver.", SLATE_300, False, 12),
        ("  Formato: moneda colombiana (COP)", SLATE_300, False, 12),
        ("  Si este número es $0 = problema grave.", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 5, 1.0, 4.5, 3, [
        ("KPI 2: Transacciones", BLUE, True, 14),
        ("  Número total de ventas aprobadas.", SLATE_300, False, 12),
        ("  = COUNT(*) de silver.ventas_clean", SLATE_300, False, 12),
        ("  Si este número es menor que Bronze", SLATE_300, False, 12),
        ("  = algunos registros fueron rechazados.", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 0.5, 3.5, 4.3, 3, [
        ("KPI 3: Calidad Promedio", AMBER, True, 14),
        ("  Promedio de quality_score de todos", SLATE_300, False, 12),
        ("  los registros de Silver.", SLATE_300, False, 12),
        ("  > 90% = saludable", GREEN, False, 12),
        ("  < 80% = necesita atención", RED, False, 12),
        ("  Medellín tiene 85% por vendedor vacío.", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 5, 3.5, 4.5, 3, [
        ("KPI 4: Registros Rechazados", RED, True, 14),
        ("  Diferencia: Bronze (8) - Silver (7) = 1", SLATE_300, False, 12),
        ("  Indica calidad de la fuente de datos.", SLATE_300, False, 12),
        ("  Si este número crece = el ERP tiene", SLATE_300, False, 12),
        ("  un problema que hay que escalar", SLATE_300, False, 12),
        ("  al Data Owner del dominio.", SLATE_300, False, 12),
    ], spacing=3)
    add_text(s, 0.5, 6.5, 9, 0.4, "Las tablas debajo muestran el detalle por ciudad y por producto con el badge de calidad.", 12, SLATE_400)

    # --- MÓDULO 8: CATÁLOGO ---
    section(prs, 8, "Catálogo de Datos", "Tab 'Catálogo' — sin catálogo, tu Lake es un pantano")

    s = slide(prs, "Tab 'Catálogo': el registro central de metadatos", "CATÁLOGO", BLUE)
    add_text(s, 0.5, 1.0, 9, 0.5, "El catálogo responde: '¿Qué datos tenemos? ¿De quién son? ¿Son confiables? ¿Tienen datos personales?'", 14, SLATE_300)
    add_text(s, 0.5, 1.7, 9, 0.3, "Cada dataset tiene estos metadatos obligatorios:", 14, BLUE, True)

    add_lines(s, 0.5, 2.1, 4.3, 4, [
        ("Identificación", WHITE, True, 13),
        ("  schema_name: bronze/silver/gold", SLATE_300, False, 11),
        ("  table_name: ventas_raw, etc.", SLATE_300, False, 11),
        ("  description: texto legible", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Responsabilidad", WHITE, True, 13),
        ("  owner: equipo-datos, equipo-iot,", SLATE_300, False, 11),
        ("  equipo-bi (Data Stewardship)", SLATE_300, False, 11),
    ], spacing=2)

    add_lines(s, 5, 2.1, 4.5, 4, [
        ("Seguridad", WHITE, True, 13),
        ("  classification: public/internal/", SLATE_300, False, 11),
        ("  confidential/restricted", SLATE_300, False, 11),
        ("  has_pii: ¿tiene datos personales?", SLATE_300, False, 11),
        ("", WHITE, False, 4),
        ("Calidad y frescura", WHITE, True, 13),
        ("  quality_score: 0.00 a 1.00", SLATE_300, False, 11),
        ("  freshness: daily/hourly/real-time", SLATE_300, False, 11),
        ("  row_count: filas actuales", SLATE_300, False, 11),
    ], spacing=2)

    add_text(s, 0.5, 5.2, 9, 0.4, "Ejemplo visual en la plataforma:", 14, AMBER, True)
    add_bullets(s, 0.5, 5.6, 9, 1.5, [
        "bronze.ventas_raw: badge BRONZE ámbar + badge CONFIDENTIAL rojo + badge PII rojo + barra 85%",
        "gold.estado_sensores: badge GOLD verde + badge PUBLIC gris + barra 97% verde",
        "La barra visual cambia de color: verde (>=95%), azul (>=90%), ámbar (>=85%), rojo (<85%)",
    ], 11)

    # --- MÓDULO 9: CALIDAD ---
    section(prs, 9, "Calidad de Datos", "Tab 'Calidad' — las 6 dimensiones del DAMA-DMBOK2")

    s = slide(prs, "Las 6 dimensiones de calidad de datos", "DAMA-DMBOK2", AMBER)
    add_text(s, 0.5, 1.0, 9, 0.4, "El DAMA-DMBOK2 define 6 dimensiones. TODAS están implementadas en la plataforma:", 14, SLATE_300)

    dims = [
        ("1. Completitud", "¿Faltan valores?", "cantidad NOT NULL", "1 registro rechazado por NULL", AMBER),
        ("2. Validez", "¿Los valores son posibles?", "cantidad > 0, precio > 0", "Caudal -5.0 L/s = anomalía", AMBER),
        ("3. Consistencia", "¿Mismo formato?", "INITCAP(ciudad)", '"medellin" → "Medellin"', AMBER),
        ("4. Unicidad", "¿Hay duplicados?", "Deduplicación por ID", "Sin duplicados en Silver", AMBER),
        ("5. Oportunidad", "¿Datos frescos?", "freshness en catálogo", "daily, hourly, real-time", AMBER),
        ("6. Confiabilidad", "¿Puedo confiar?", "quality_score por registro", "0.70 = sospechoso, 1.00 = ok", AMBER),
    ]

    for i, (dim, pregunta, regla, ejemplo, color) in enumerate(dims):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.7
        y = 1.6 + row * 1.6
        add_text(s, x, y, 4.3, 0.3, dim, 13, WHITE, True)
        add_lines(s, x, y + 0.3, 4.3, 1, [
            (f"  Pregunta: {pregunta}", SLATE_300, False, 10),
            (f"  Regla: {regla}", CYAN, False, 10),
            (f"  Ejemplo: {ejemplo}", SLATE_400, False, 10),
        ], spacing=1)

    add_text(s, 0.5, 6.5, 9, 0.4, "Tab 'Calidad' muestra: 5 métricas arriba + 7 reglas activas + problemas detectados + barras por dataset.", 11, GREEN)

    # Calidad en la plataforma
    s = slide(prs, "Tab 'Calidad': cómo interpretarla paso a paso", "OPERACIÓN", AMBER)
    add_text(s, 0.5, 1.0, 9, 0.3, "Al abrir esta tab, lee de arriba hacia abajo:", 14, SLATE_300)

    add_lines(s, 0.5, 1.5, 9, 5.5, [
        ("Sección 1 — Métricas globales (5 tarjetas)", WHITE, True, 14),
        ("  Registros Bronze: 8 | Aprobados Silver: 7 | Rechazados: 1 | Tasa: 87% | Score: 96%", SLATE_300, False, 12),
        ("  Si la tasa baja = el sistema fuente está enviando peores datos. Escalar al Data Owner.", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("Sección 2 — Barra visual del pipeline", WHITE, True, 14),
        ("  Barra verde (7 aprobados) + roja (1 rechazado) = visualización inmediata del filtro.", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Sección 3 — Reglas de calidad activas (tabla izquierda)", WHITE, True, 14),
        ("  7 reglas con su tipo (Completitud, Validez, Consistencia, Unicidad, Confiabilidad).", SLATE_300, False, 12),
        ("  Cada regla tiene estado: 'Activa — filtró registros' o 'Activa — normaliza'.", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Sección 4 — Problemas detectados (panel derecho)", WHITE, True, 14),
        ("  3 problemas en Bronze, cada uno con severidad:", SLATE_300, False, 12),
        ("  🔴 CRÍTICO: cantidad NULL (bloquea) | 🟡 WARNING: vendedor vacío (corrige) | 🔵 INFO: ciudad (normaliza)", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Sección 5 — Quality Score por dataset (barras abajo)", WHITE, True, 14),
        ("  Cada dataset del catálogo con su barra visual de calidad. Verde > Azul > Ámbar > Rojo.", SLATE_300, False, 12),
    ], spacing=1)

    # --- MÓDULO 10: LINAJE ---
    section(prs, 10, "Linaje de Datos (Data Lineage)", "Tab 'Linaje' — ¿de dónde viene cada dato?")

    s = slide(prs, "¿Qué es el linaje y por qué es crítico?", "LINAJE", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.8, "El linaje de datos (data lineage) documenta el ORIGEN, las TRANSFORMACIONES y el DESTINO de cada dato. Responde: 'este KPI de $22M en Barranquilla... ¿de dónde sale?'", 14, SLATE_300)
    add_text(s, 0.5, 2.0, 9, 0.4, "¿Cuándo necesitas el linaje?", 14, PURPLE, True)
    add_bullets(s, 0.5, 2.4, 9, 2, [
        "Cuando un KPI parece incorrecto y necesitas rastrear la causa",
        "Cuando un regulador pregunta cómo se calculó un reporte",
        "Cuando cambias un pipeline y necesitas saber qué reportes se afectan (impact analysis)",
        "Cuando necesitas cumplir GDPR Art. 15: derecho de acceso — explicar qué datos tienes y de dónde vienen",
    ], 12)
    add_text(s, 0.5, 4.2, 9, 0.4, "Tab 'Linaje' en la plataforma tiene 2 secciones:", 14, GREEN, True)
    add_lines(s, 0.5, 4.6, 9, 2.5, [
        ("1. Grafo macro: FUENTE → BRONZE → SILVER → GOLD → CONSUMO", WHITE, True, 13),
        ("   Cada nivel muestra las tablas, transformaciones aplicadas y conteo de registros.", SLATE_300, False, 11),
        ("   SAP ERP + IoT Gateway → ventas_raw + sensores_raw → ventas_clean → ventas_por_ciudad → Dashboard", SLATE_400, False, 11),
        ("", WHITE, False, 4),
        ("2. Tabla de linaje por registro:", WHITE, True, 13),
        ("   Silver ID → Producto → Ciudad → Calidad → Bronze ID → Sistema → Archivo", SLATE_300, False, 11),
        ("   S#1 → Sensor IoT X200 → Medellín → 100% ← B#1 ← SAP-ERP ← export_ventas_20260301.json", SLATE_400, False, 11),
    ], spacing=2)

    # Linaje SQL
    s = slide(prs, "Linaje: la implementación técnica", "SQL", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.4, "El linaje se implementa con una foreign key de Silver a Bronze:", 14, SLATE_300)
    add_code(s, 0.5, 1.5, 9, 1.8,
"""-- En silver.ventas_clean:
source_record_id BIGINT REFERENCES bronze.ventas_raw(id)
-- Cada registro Silver sabe exactamente de qué registro Bronze vino.

-- Para rastrear un dato sospechoso en Gold:
-- Paso 1: ¿Qué registros Silver componen el KPI de Barranquilla?
SELECT * FROM silver.ventas_clean WHERE ciudad = 'Barranquilla';
-- Paso 2: ¿De qué registro Bronze vino?
SELECT * FROM bronze.ventas_raw WHERE id = 5;  -- source_record_id = 5
-- Resultado: Bomba Sumergible BS30, 25 uds × $890K = $22.250.000 ✓""", 10)

    add_text(s, 0.5, 3.6, 9, 0.4, "En producción, herramientas que automatizan esto:", 14, AMBER, True)
    add_lines(s, 0.5, 4.0, 4.3, 3, [
        ("Apache Atlas", WHITE, True, 13),
        ("  Open source (Apache Foundation)", SLATE_300, False, 11),
        ("  Integrado con Hadoop/Hive/Spark", SLATE_300, False, 11),
        ("", WHITE, False, 3),
        ("OpenLineage", WHITE, True, 13),
        ("  Estándar abierto para linaje", SLATE_300, False, 11),
        ("  Integrado con Airflow, dbt, Spark", SLATE_300, False, 11),
    ], spacing=2)
    add_lines(s, 5, 4.0, 4.5, 3, [
        ("dbt lineage", WHITE, True, 13),
        ("  Linaje automático desde SQL", SLATE_300, False, 11),
        ("  Grafo visual en dbt Cloud", SLATE_300, False, 11),
        ("", WHITE, False, 3),
        ("Microsoft Purview", WHITE, True, 13),
        ("  Enterprise governance platform", SLATE_300, False, 11),
        ("  Catálogo + linaje + clasificación", SLATE_300, False, 11),
    ], spacing=2)

    # ================================================================
    # PARTE III: SEGURIDAD Y OPERACIÓN
    # ================================================================

    # --- MÓDULO 11: SEGURIDAD ---
    section(prs, 11, "Seguridad de Datos", "Tabs 'Usuarios' — RLS, enmascaramiento y RBAC")

    s = slide(prs, "Los 4 roles y sus permisos exactos", "RBAC", CYAN)
    add_text(s, 0.5, 1.0, 9, 0.4, "Principio de mínimo privilegio: cada usuario accede SOLO a lo que necesita para su trabajo.", 14, SLATE_300)

    add_lines(s, 0.5, 1.6, 4.3, 2.5, [
        ("ADMIN — Admin Lake (TI)", RED, True, 14),
        ("  Bronze R/W + Silver R/W + Gold R/W", SLATE_300, False, 12),
        ("  Audit R + Admin: gestiona todo", SLATE_300, False, 12),
        ("  Solo personal de TI/Gobernanza", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 5, 1.6, 4.5, 2.5, [
        ("DATA ENGINEER — Carlos (TI)", GREEN, True, 14),
        ("  Bronze R/W + Silver R/W + Gold R", SLATE_300, False, 12),
        ("  Carga datos, construye pipelines", SLATE_300, False, 12),
        ("  No puede modificar Gold ni Audit", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 0.5, 3.8, 4.3, 2.5, [
        ("DATA SCIENTIST — María (Innov.)", PURPLE, True, 14),
        ("  Silver MASKED R + Gold R + Sandbox", SLATE_300, False, 12),
        ("  Ve ventas_masked, NO ventas_clean", SLATE_300, False, 12),
        ("  NIT: '900****56' (no ve el real)", SLATE_300, False, 12),
    ], spacing=3)
    add_lines(s, 5, 3.8, 4.5, 2.5, [
        ("ANALYST — Pedro (Comercial)", BLUE, True, 14),
        ("  Gold R solamente", SLATE_300, False, 12),
        ("  Solo datos agregados y limpios", SLATE_300, False, 12),
        ("  No ve Bronze, Silver, ni PII", SLATE_300, False, 12),
    ], spacing=3)

    add_text(s, 0.5, 6.2, 9, 0.5, "En la tab 'Usuarios': cada tarjeta muestra el nombre, email, rol con color, departamento y permisos exactos.", 12, SLATE_400)

    # PII Masking
    s = slide(prs, "Enmascaramiento de PII: proteger datos personales", "PII", RED)
    add_text(s, 0.5, 1.0, 9, 0.5, "PII = Personally Identifiable Information. En Colombia, la Ley 1581 exige protección especial.", 14, SLATE_300)
    add_text(s, 0.5, 1.7, 9, 0.4, "Datos PII en nuestra plataforma:", 14, RED, True)
    add_bullets(s, 0.5, 2.1, 9, 1, [
        "cliente_nit: número de identificación tributaria (NIT) — identifica al cliente",
        "cliente_nombre: nombre de la empresa — puede identificar personas naturales",
    ], 12)
    add_text(s, 0.5, 3.0, 9, 0.4, "La vista enmascarada:", 14, GREEN, True)
    add_code(s, 0.5, 3.4, 9, 2.5,
"""-- silver.ventas_masked — lo que ve María (data_scientist):
CREATE VIEW silver.ventas_masked AS
SELECT
    id, fecha, producto, cantidad, precio_unit, vendedor, ciudad,

    -- NIT: "900123456" → "900****56"
    CONCAT(LEFT(cliente_nit, 3), '****', RIGHT(cliente_nit, 2))
        AS cliente_nit_masked,

    -- Nombre: "Acueducto Municipal" → "Acu***"
    CONCAT(LEFT(cliente_nombre, 3), '***')
        AS cliente_nombre_masked,

    ingreso_total, quality_score
FROM silver.ventas_clean;""", 10)
    add_text(s, 0.5, 6.0, 9, 0.8, "María puede analizar patrones de ventas por producto y ciudad sin nunca ver la identidad real del cliente. Esto cumple con el principio de minimización de datos del GDPR y Habeas Data.", 12, SLATE_400)

    # RLS
    s = slide(prs, "Row Level Security (RLS): control a nivel de fila", "RLS", RED)
    add_text(s, 0.5, 1.0, 9, 0.4, "RLS es una feature nativa de PostgreSQL que filtra filas según el usuario conectado.", 14, SLATE_300)
    add_code(s, 0.5, 1.5, 9, 2.5,
"""-- 1. Habilitar RLS (sin esto, todos ven todo)
ALTER TABLE silver.ventas_clean ENABLE ROW LEVEL SECURITY;

-- 2. Política: solo engineers y admins ven Silver completo
CREATE POLICY "engineers_full_access"
    ON silver.ventas_clean
    FOR ALL
    USING (
        EXISTS (
            SELECT 1 FROM public.lake_users
            WHERE email = current_setting('app.current_user_email')
            AND role IN ('data_engineer', 'admin')
        )
    );

-- Sin esta política, un SELECT a silver.ventas_clean devuelve 0 filas
-- para cualquier usuario que no sea engineer o admin.""", 10)
    add_text(s, 0.5, 4.3, 9, 0.4, "Diferencia práctica:", 14, AMBER, True)
    add_lines(s, 0.5, 4.7, 4.3, 2, [
        ("Sin RLS:", RED, True, 13),
        ("  SELECT * FROM silver.ventas_clean", SLATE_300, False, 11),
        ("  → 7 filas (todos ven todo)", SLATE_300, False, 11),
        ("  Riesgo: NIT visible para analistas", RED, False, 11),
    ], spacing=3)
    add_lines(s, 5, 4.7, 4.5, 2, [
        ("Con RLS:", GREEN, True, 13),
        ("  Pedro (analyst) → 0 filas", SLATE_300, False, 11),
        ("  María → vía ventas_masked (PII oculto)", SLATE_300, False, 11),
        ("  Carlos (engineer) → 7 filas completas", GREEN, False, 11),
    ], spacing=3)

    # --- MÓDULO 12: AUDITORÍA ---
    section(prs, 12, "Auditoría de Accesos", "Tab 'Auditoría' — cumplimiento regulatorio")

    s = slide(prs, "Tab 'Auditoría': cada acceso queda registrado", "AUDITORÍA", RED)
    add_text(s, 0.5, 1.0, 9, 0.4, "La auditoría responde: '¿Quién accedió a qué datos, cuándo, cuánto y desde dónde?'", 14, SLATE_300)
    add_text(s, 0.5, 1.5, 9, 0.4, "Columnas del log:", 14, RED, True)
    add_lines(s, 0.5, 1.9, 9, 2.5, [
        ("  timestamp: fecha y hora exacta del acceso (ej: 15/03/2026 10:30:22)", SLATE_300, False, 12),
        ("  user_email: quién (ej: pedro.analista@eafit.edu.co)", SLATE_300, False, 12),
        ("  user_role: rol del usuario (ej: analyst, data_engineer)", SLATE_300, False, 12),
        ("  action: qué hizo → SELECT (consultó), INSERT (cargó), EXPORT (descargó)", SLATE_300, False, 12),
        ("  schema.table: qué tabla (ej: gold.ventas_por_ciudad)", SLATE_300, False, 12),
        ("  row_count: cuántas filas accedió (ej: 5)", SLATE_300, False, 12),
        ("  ip_address: desde qué red (ej: 10.0.3.8)", SLATE_300, False, 12),
    ], spacing=2)
    add_text(s, 0.5, 4.7, 9, 0.4, "Escenarios reales donde usarías este log:", 14, AMBER, True)
    add_bullets(s, 0.5, 5.1, 9, 2, [
        "Regulador pregunta: '¿Quién accedió a datos PII en el último mes?' → filtrar por tabla + PII",
        "Incidente: '¿Quién descargó (EXPORT) datos de clientes?' → filtrar action='EXPORT'",
        "Verificación: '¿Pedro (analyst) intentó acceder a Bronze?' → filtrar email + schema='bronze'",
        "Compliance: '¿Todos los accesos a datos confidenciales fueron de personal autorizado?'",
    ], 12)

    # --- MÓDULO 13: OPERACIÓN ---
    section(prs, 13, "Operando la plataforma", "Guía paso a paso por las 11 tabs")

    s = slide(prs, "Sensores IoT: detección de anomalías en Gold", "IOT", GREEN)
    add_text(s, 0.5, 1.0, 9, 0.4, "La tab 'Sensores IoT' consume gold.estado_sensores — una vista materializada.", 14, SLATE_300)
    add_text(s, 0.5, 1.5, 9, 0.4, "Cada tarjeta muestra:", 14, GREEN, True)
    add_bullets(s, 0.5, 1.9, 9, 1.5, [
        "Sensor ID + ubicación + tipo (caudal, presión, calidad, nivel)",
        "4 métricas: promedio, mínimo, máximo, total lecturas",
        "Batería con color: verde >70%, ámbar 50-70%, rojo <50% (SEN-004 = 45%)",
        "Badge rojo 'Anomalía' si alguna lectura es físicamente imposible",
    ], 12)
    add_text(s, 0.5, 3.5, 9, 0.4, "Caso de anomalía — SEN-001 (Planta Norte):", 14, RED, True)
    add_code(s, 0.5, 3.9, 9, 1.8,
"""-- SEN-001 tiene 2 lecturas de caudal:
-- Lectura 1: 245.7 L/s (normal — caudal de planta)
-- Lectura 2: -5.0 L/s  (IMPOSIBLE — caudal no puede ser negativo)

-- Pipeline Bronze→Silver marcó la lectura negativa:
is_anomaly = CASE WHEN valor < 0 THEN TRUE ELSE FALSE END

-- Gold agrega con BOOL_OR:
tiene_anomalias = BOOL_OR(is_anomaly)  -- TRUE si CUALQUIER lectura es anómala

-- En la plataforma: SEN-001 muestra badge rojo "Anomalía detectada"
-- Acción: el operador debe revisar el sensor (calibración o falla)""", 10)
    add_text(s, 0.5, 6.0, 9, 0.5, "En producción: estas anomalías dispararían alertas por Slack/PagerDuty al equipo de operaciones.", 12, SLATE_400)

    # Flujo de navegación
    s = slide(prs, "Flujo recomendado para explorar la plataforma", "GUÍA", BLUE)
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("Paso 1 → Tab 'Bronze (Raw)'", AMBER, True, 15),
        ("    Ve los 8 registros JSONB. Identifica los 3 errores del registro #4.", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 2 → Tab 'Bronze vs Silver'", CYAN, True, 15),
        ("    Compara antes/después. Encuentra el registro RECHAZADO y los TRANSFORMADOS.", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 3 → Tab 'Silver (Clean)'", PURPLE, True, 15),
        ("    Ve los 7 registros limpios. Nota quality_score, ingreso_total calculado, source_record_id.", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 4 → Tab 'Dashboard Gold'", GREEN, True, 15),
        ("    Lee los 4 KPIs. ¿Por qué Medellín tiene 85% de calidad? ¿Cuántos rechazados hay?", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 5 → Tab 'Calidad'", AMBER, True, 15),
        ("    Revisa las 7 reglas activas. ¿Qué problemas se detectaron? ¿Cuál es la tasa de aprobación?", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 6 → Tab 'Linaje'", PURPLE, True, 15),
        ("    Sigue el grafo FUENTE→BRONZE→SILVER→GOLD. Rastrea un registro Silver hasta su Bronze.", SLATE_300, False, 12),
        ("", WHITE, False, 3),
        ("Paso 7 → Tabs 'Catálogo' + 'Auditoría' + 'Usuarios' + 'Sensores' + 'Arquitectura'", WHITE, True, 15),
        ("    Explora metadatos, accesos, roles, anomalías IoT y el stack técnico completo.", SLATE_300, False, 12),
    ], spacing=1)

    # --- MÓDULO 14: MADUREZ ---
    section(prs, 14, "Modelo de Madurez y Siguientes Pasos", "¿Dónde estamos y hacia dónde vamos?")

    s = slide(prs, "Modelo de Madurez de Gobernanza de Datos", "MADUREZ", PURPLE)
    add_text(s, 0.5, 1.0, 9, 0.4, "Basado en CMMI y adaptado a Data Governance (Stanford/DAMA):", 13, SLATE_300)

    levels = [
        ("Nivel 1 — Inicial (Ad-hoc)", "No hay procesos. Cada equipo maneja sus datos.", SLATE_400),
        ("Nivel 2 — Repetible", "Procesos básicos. Algunos estándares. Sin automatización.", AMBER),
        ("Nivel 3 — Definido ← NUESTRA PLATAFORMA", "Procesos documentados. Catálogo, calidad, linaje, roles definidos.", GREEN),
        ("Nivel 4 — Gestionado", "Métricas de calidad monitoreadas. SLAs definidos. Alertas automáticas.", BLUE),
        ("Nivel 5 — Optimizado", "Mejora continua. ML para calidad. Self-service data. DataOps.", PURPLE),
    ]
    for i, (level, desc, color) in enumerate(levels):
        y = 1.6 + i * 1.0
        add_text(s, 0.5, y, 9, 0.3, level, 13, color, True)
        add_text(s, 0.5, y + 0.3, 9, 0.3, f"    {desc}", 11, SLATE_300)

    add_text(s, 0.5, 6.6, 9, 0.4, "Nuestra plataforma está en Nivel 3: procesos definidos, catálogo, calidad, linaje, roles. Para Nivel 4, necesitamos alertas automáticas y SLAs.", 11, SLATE_400)

    # Siguientes pasos
    s = slide(prs, "Siguientes pasos para llevar esto a producción", "FUTURO", GREEN)
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("Evolución técnica", BLUE, True, 15),
        ("  Migrar a un Data Lake real: S3 (AWS), ADLS (Azure) o GCS (GCP)", SLATE_300, False, 12),
        ("  Usar Parquet en vez de PostgreSQL para volúmenes masivos (TB+)", SLATE_300, False, 12),
        ("  Implementar Delta Lake o Apache Iceberg para ACID + Time Travel", SLATE_300, False, 12),
        ("  Pipeline con Airflow/dbt en vez de SQL manual", SLATE_300, False, 12),
        ("  Catálogo con Apache Atlas o Microsoft Purview", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Evolución de gobernanza", PURPLE, True, 15),
        ("  Definir Data Owners y Stewards formales por dominio", SLATE_300, False, 12),
        ("  Crear comité de gobernanza con reuniones periódicas", SLATE_300, False, 12),
        ("  Implementar SLAs de calidad: 'ventas_clean debe tener calidad >= 95%'", SLATE_300, False, 12),
        ("  Automatizar alertas cuando quality_score baje de umbral", SLATE_300, False, 12),
        ("  Glosario de negocio: definir términos como 'venta activa', 'cliente recurrente'", SLATE_300, False, 12),
        ("", WHITE, False, 4),
        ("Evolución regulatoria", RED, True, 15),
        ("  Implementar Data Protection Impact Assessment (DPIA) para nuevos datasets", SLATE_300, False, 12),
        ("  Proceso formal de solicitud de acceso a datos (workflow de aprobación)", SLATE_300, False, 12),
        ("  Retención automatizada: borrar datos Bronze después de X años (lifecycle policies)", SLATE_300, False, 12),
    ], spacing=1)

    # Preguntas de comprensión
    s = slide(prs, "Preguntas de comprensión", "REFLEXIÓN", AMBER)
    add_lines(s, 0.5, 1.0, 9, 6.2, [
        ("1. ¿Por qué el registro #4 fue rechazado de Silver?", WHITE, True, 13),
        ("   Porque tiene cantidad = NULL. La regla WHERE cantidad IS NOT NULL lo filtra.", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("2. ¿Qué diferencia hay entre Gobernanza y Gestión de datos?", WHITE, True, 13),
        ("   Gobernanza DEFINE las políticas (ej: 'PII debe enmascararse'). Gestión las IMPLEMENTA.", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("3. ¿Por qué María (data_scientist) no ve el NIT real?", WHITE, True, 13),
        ("   Porque accede vía ventas_masked, que aplica CONCAT(LEFT(nit,3),'****',RIGHT(nit,2)).", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("4. ¿Qué regulación colombiana exige proteger datos personales?", WHITE, True, 13),
        ("   Ley 1581 de 2012 (Habeas Data). Exige consentimiento, finalidad y acceso controlado.", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("5. ¿Qué dimensión de calidad DAMA mide si hay valores faltantes?", WHITE, True, 13),
        ("   Completitud. Regla: cantidad NOT NULL. 1 registro rechazado por fallar esta dimensión.", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("6. ¿Cómo rastrearías un KPI sospechoso de Gold hasta su origen?", WHITE, True, 13),
        ("   Gold (agregación) → Silver (source_record_id) → Bronze (JSONB original del ERP).", SLATE_400, False, 11),
        ("", WHITE, False, 3),
        ("7. ¿Qué pasó con el sensor SEN-001 y por qué importa para gobernanza?", WHITE, True, 13),
        ("   Midió -5.0 L/s (imposible). El pipeline lo marcó como anomalía. La gobernanza exige", SLATE_400, False, 11),
        ("   detectar y alertar sobre datos que no tienen sentido físico.", SLATE_400, False, 11),
    ], spacing=1)

    # Resumen conceptos
    s = slide(prs, "Glosario de la lección", "REFERENCIA", BLUE)
    add_lines(s, 0.5, 1.0, 4.3, 6.2, [
        ("DAMA-DMBOK2", WHITE, True, 12),
        ("  Framework de gestión de datos (11 áreas)", SLATE_400, False, 10),
        ("Data Lake", WHITE, True, 12),
        ("  Repositorio centralizado, cualquier formato", SLATE_400, False, 10),
        ("Medallion", WHITE, True, 12),
        ("  Arquitectura Bronze → Silver → Gold", SLATE_400, False, 10),
        ("Schema-on-read", WHITE, True, 12),
        ("  Esquema se define al consultar, no al cargar", SLATE_400, False, 10),
        ("ETL / ELT", WHITE, True, 12),
        ("  Extract-Transform-Load / Extract-Load-Transform", SLATE_400, False, 10),
        ("Quality Score", WHITE, True, 12),
        ("  Métrica numérica de confiabilidad por registro", SLATE_400, False, 10),
        ("Data Lineage", WHITE, True, 12),
        ("  Trazabilidad origen → transformación → destino", SLATE_400, False, 10),
        ("RLS", WHITE, True, 12),
        ("  Row Level Security — control de acceso por fila", SLATE_400, False, 10),
    ], spacing=2)
    add_lines(s, 5, 1.0, 4.5, 6.2, [
        ("PII", WHITE, True, 12),
        ("  Personally Identifiable Information", SLATE_400, False, 10),
        ("RBAC", WHITE, True, 12),
        ("  Role-Based Access Control", SLATE_400, False, 10),
        ("PII Masking", WHITE, True, 12),
        ("  Ocultar datos personales en vistas", SLATE_400, False, 10),
        ("Habeas Data", WHITE, True, 12),
        ("  Ley 1581/2012 — protección datos Colombia", SLATE_400, False, 10),
        ("GDPR", WHITE, True, 12),
        ("  General Data Protection Regulation (EU)", SLATE_400, False, 10),
        ("Materialized View", WHITE, True, 12),
        ("  Consulta SQL pre-calculada como tabla", SLATE_400, False, 10),
        ("Data Steward", WHITE, True, 12),
        ("  Custodio operativo de un dominio de datos", SLATE_400, False, 10),
        ("Data Owner", WHITE, True, 12),
        ("  Responsable estratégico de un dominio", SLATE_400, False, 10),
    ], spacing=2)

    # CIERRE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, SLATE_900)
    add_text(s, 0.5, 1.5, 9, 1, "Un Data Lake sin gobernanza\nes solo un bucket con datos.", 34, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, 0.5, 3.2, 9, 0.8, "datalake-governance-hub.vercel.app", 22, BLUE, False, PP_ALIGN.CENTER)
    add_text(s, 0.5, 4.3, 9, 1.5, "Navega las 11 tabs. Compara Bronze vs Silver.\nRastrea el linaje. Revisa la auditoría.\nEntiende los roles. Explora el catálogo.\n\nLa gobernanza de datos no es un proyecto,\nes una práctica continua.", 15, SLATE_300, False, PP_ALIGN.CENTER)
    add_text(s, 0.5, 6.5, 9, 0.5, "Gestión y Gobernanza de Datos — Universidad EAFIT — 2026", 12, SLATE_500, False, PP_ALIGN.CENTER)

    # SAVE
    path = os.path.expanduser("~/Documents/datalake-governance-hub/PPTX_2_Como_Funciona_Plataforma.pptx")
    prs.save(path)
    print(f"✅ Guardado: {path}")
    print(f"   Total slides: {len(prs.slides)}")
    return path

if __name__ == '__main__':
    build()
