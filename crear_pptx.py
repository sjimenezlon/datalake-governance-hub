#!/usr/bin/env python3
"""
Genera dos PowerPoints para el curso de Gestión y Gobernanza de Datos EAFIT:
1. Paso a paso para replicar la actividad
2. Cómo funciona la plataforma con ejemplos
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== COLORES =====
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
SLATE_900 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_800 = RGBColor(0x1E, 0x40, 0x5F)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_DARK = RGBColor(0x1A, 0x56, 0xDB)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=SLATE_300, bullet_color=BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_numbered_list(slide, left, top, width, height, items, font_size=14, color=SLATE_300):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox

def add_code_box(slide, left, top, width, height, code, font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p.font.name = 'Consolas'
    return shape

def add_badge(slide, left, top, text, bg_color, text_color=WHITE, width=1.2, height=0.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 2.0, 8.4, 1.5, title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 3.5, 8.4, 1.0, subtitle, font_size=16, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    # Footer
    add_text_box(slide, 0.8, 6.5, 8.4, 0.5, "Gestión y Gobernanza de Datos — Universidad EAFIT — 2026", font_size=10, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(prs, section_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLATE_900)
    # Section number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.3), Inches(1.8), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, 0.8, 2.8, 8.4, 1.0, title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 0.8, 3.8, 8.4, 0.8, subtitle, font_size=14, color=SLATE_400, alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(prs, title, badge_text=None, badge_color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.5, 0.3, 7, 0.6, title, font_size=24, bold=True, color=WHITE)
    if badge_text:
        add_badge(slide, 7.8, 0.35, badge_text, badge_color)
    # Separator line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(9), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SLATE_800
    shape.line.fill.background()
    return slide


# ================================================================
# PPTX 1: PASO A PASO PARA REPLICAR LA ACTIVIDAD
# ================================================================
def create_paso_a_paso():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: Portada ---
    title_slide(prs,
        "Data Lake Governance Hub\nGuía Paso a Paso",
        "Cómo replicar el laboratorio completo\nSupabase + Next.js + Vercel")

    # --- SLIDE 2: Qué vamos a construir ---
    slide = content_slide(prs, "¿Qué vamos a construir?", "OVERVIEW", PURPLE)
    add_bullet_list(slide, 0.5, 1.2, 9, 5.5, [
        "Un Data Lake con arquitectura Medallion (Bronze → Silver → Gold) usando PostgreSQL en Supabase",
        "Datos reales de ventas (SAP ERP simulado) y sensores IoT con errores intencionales",
        "Transformación automática Bronze→Silver: limpieza, validación, normalización",
        "Vistas materializadas Gold para KPIs de negocio",
        "Row Level Security (RLS) — control de acceso por rol",
        "Enmascaramiento de datos PII (NIT, nombre de cliente)",
        "Catálogo de datos con metadatos obligatorios",
        "Auditoría de accesos — quién consultó qué, cuándo y desde dónde",
        "Dashboard web desplegado en Vercel que SOLO consume capa Gold",
        "11 tabs interactivas: Dashboard, Comparación, Calidad, Linaje, Bronze, Silver, Catálogo, Sensores, Auditoría, Usuarios, Arquitectura",
    ], font_size=13)

    # --- SLIDE 3: Requisitos previos ---
    slide = content_slide(prs, "Requisitos previos", "PASO 0", AMBER)
    add_numbered_list(slide, 0.5, 1.2, 9, 5.5, [
        "Cuenta de GitHub (github.com) — gratuita",
        "Cuenta de Supabase (supabase.com) — plan gratuito, login con GitHub",
        "Cuenta de Vercel (vercel.com) — plan gratuito, login con GitHub",
        "Node.js instalado (versión 18+) — nodejs.org",
        "Editor de código: VS Code recomendado (code.visualstudio.com)",
        "Terminal / línea de comandos (CMD, PowerShell, Terminal de Mac)",
    ], font_size=15)
    add_text_box(slide, 0.5, 5.5, 9, 1, "Tiempo estimado total: ~90 minutos", font_size=14, color=AMBER, bold=True)

    # --- SLIDE 4: Crear proyecto Supabase ---
    section_slide(prs, 1, "Crear proyecto en Supabase", "Base de datos PostgreSQL en la nube")

    slide = content_slide(prs, "Paso 1: Crear proyecto Supabase", "SUPABASE", GREEN)
    add_numbered_list(slide, 0.5, 1.2, 9, 3.5, [
        'Ir a supabase.com → "Start your project" → Login con GitHub',
        'Clic en "New Project"',
        'Nombre: datalake-gobernanza-eafit',
        'Región: South America (São Paulo) — la más cercana a Colombia',
        'Database Password: genera una segura y guárdala',
        'Clic en "Create new project" — esperar ~2 minutos',
    ], font_size=14)
    add_text_box(slide, 0.5, 5.0, 9, 0.5, "Una vez creado, ve a Settings → API y copia:", font_size=13, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 5.5, 9, 1.5, [
        "Project URL → algo como https://xyzabc123.supabase.co",
        "anon public key → string largo que empieza con eyJ...",
    ], font_size=13, color=SLATE_300)

    # --- SLIDE 5: Ejecutar SQL ---
    section_slide(prs, 2, "Ejecutar SQL para crear el Data Lake", "Esquemas, tablas, datos, vistas y permisos")

    slide = content_slide(prs, "Paso 2: Ejecutar SQL en Supabase", "SQL", BLUE)
    add_numbered_list(slide, 0.5, 1.2, 9, 2.5, [
        "En el menú izquierdo de Supabase → SQL Editor",
        'Clic en "New query"',
        "Abrir el archivo supabase_setup.sql (incluido en el repo)",
        "Seleccionar todo (Ctrl+A / Cmd+A), copiar (Ctrl+C / Cmd+C)",
        "Pegar en el SQL Editor de Supabase (Ctrl+V / Cmd+V)",
        'Clic en "Run" — debe decir "Success. No rows returned"',
    ], font_size=14)
    add_text_box(slide, 0.5, 4.5, 9, 0.5, "Este SQL crea:", font_size=13, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 5.0, 4.5, 2.5, [
        "4 esquemas: bronze, silver, gold, audit",
        "8 tablas con datos de prueba",
        "3 vistas materializadas Gold",
        "9 vistas públicas para la API",
    ], font_size=12)
    add_bullet_list(slide, 5, 5.0, 4.5, 2.5, [
        "Permisos de lectura para anon/authenticated",
        "Usuarios con 4 roles diferentes",
        "Datos con errores intencionales",
        "Auditoría de accesos simulada",
    ], font_size=12)

    # --- SLIDE 6: Qué hace el SQL internamente ---
    slide = content_slide(prs, "¿Qué hace el SQL? — Transformación Bronze → Silver", "ETL", AMBER)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "El INSERT INTO silver.ventas_clean hace estas transformaciones:", font_size=14, color=SLATE_300)
    add_code_box(slide, 0.5, 1.8, 9, 3.5,
"""-- De Bronze (JSONB crudo) a Silver (columnas tipadas):
INSERT INTO silver.ventas_clean (...)
SELECT
  (raw_data->>'fecha')::date,         -- Casteo de texto a fecha
  (raw_data->>'cantidad')::integer,   -- Casteo a número entero
  COALESCE(NULLIF(vendedor,''), 'Sin asignar'),  -- Reemplazo de vacíos
  INITCAP(raw_data->>'ciudad'),       -- "medellin" → "Medellín"
  CASE WHEN vendedor = '' THEN 0.70 ELSE 1.00 END  -- Quality score
FROM bronze.ventas_raw
WHERE (raw_data->>'cantidad') IS NOT NULL  -- Filtra NULLs
  AND (raw_data->>'cantidad')::integer > 0; -- Filtra negativos""", font_size=11)
    add_text_box(slide, 0.5, 5.5, 9, 1.5, "Resultado: 8 registros Bronze → 7 aprobados en Silver (1 rechazado por cantidad NULL). El registro rechazado permanece en Bronze como evidencia.", font_size=13, color=SLATE_400)

    # --- SLIDE 7: Crear proyecto Next.js ---
    section_slide(prs, 3, "Crear la aplicación web", "Next.js + Tailwind CSS")

    slide = content_slide(prs, "Paso 3: Clonar el repositorio", "NEXT.JS", BLUE)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Opción A — Clonar desde GitHub (recomendado):", font_size=14, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 1.8, 9, 1.2,
"""git clone https://github.com/sjimenezlon/datalake-governance-hub.git
cd datalake-governance-hub
npm install""", font_size=13)

    add_text_box(slide, 0.5, 3.3, 9, 0.5, "Configurar las variables de entorno:", font_size=14, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 3.9, 9, 1.2,
"""# Crear archivo .env.local en la raíz del proyecto
NEXT_PUBLIC_SUPABASE_URL=https://TU-PROYECTO.supabase.co
NEXT_PUBLIC_SUPABASE_ANON_KEY=tu-anon-key-aqui""", font_size=13)

    add_text_box(slide, 0.5, 5.4, 9, 0.5, "Ejecutar en local:", font_size=14, color=GREEN, bold=True)
    add_code_box(slide, 0.5, 5.9, 9, 0.8,
"""npm run dev
# Abrir http://localhost:3000""", font_size=13)

    # --- SLIDE 8: Desplegar en Vercel ---
    section_slide(prs, 4, "Desplegar en Vercel", "Tu app en producción en 2 minutos")

    slide = content_slide(prs, "Paso 4: Desplegar en Vercel", "VERCEL", CYAN)
    add_numbered_list(slide, 0.5, 1.2, 9, 2, [
        "Ir a vercel.com → Login con GitHub",
        'Clic en "Add New Project"',
        "Importar el repositorio datalake-governance-hub de tu GitHub",
        'En "Environment Variables" agregar:',
    ], font_size=14)
    add_code_box(slide, 0.8, 3.5, 8.5, 1,
"""NEXT_PUBLIC_SUPABASE_URL = https://TU-PROYECTO.supabase.co
NEXT_PUBLIC_SUPABASE_ANON_KEY = eyJ...(tu anon key)""", font_size=12)
    add_numbered_list(slide, 0.5, 4.8, 9, 2, [
        'Clic en "Deploy" — Vercel construye y despliega automáticamente',
        "Tu app estará disponible en: https://tu-proyecto.vercel.app",
        "Cada push a main en GitHub redespliega automáticamente",
    ], font_size=14)

    # --- SLIDE 9: Estructura del proyecto ---
    slide = content_slide(prs, "Estructura del proyecto", "CÓDIGO", PURPLE)
    add_code_box(slide, 0.5, 1.2, 9, 5.5,
"""datalake-governance-hub/
├── src/
│   ├── app/
│   │   ├── layout.tsx          ← Layout global (metadatos, fuente)
│   │   ├── page.tsx            ← Página principal (11 tabs)
│   │   └── globals.css         ← Estilos Tailwind
│   ├── components/
│   │   ├── QualityTab.tsx      ← Dashboard de calidad de datos
│   │   ├── LineageTab.tsx      ← Visualización de linaje
│   │   └── CompareTab.tsx      ← Comparación Bronze vs Silver
│   └── lib/
│       └── supabase.ts         ← Cliente Supabase + tipos TypeScript
├── supabase_setup.sql          ← SQL para crear todo en Supabase
├── .env.local                  ← Variables de entorno (NO se sube a Git)
├── package.json                ← Dependencias del proyecto
├── tailwind.config.ts          ← Configuración de Tailwind CSS
└── tsconfig.json               ← Configuración de TypeScript""", font_size=12)

    # --- SLIDE 10: Checklist de entrega ---
    slide = content_slide(prs, "Checklist de entrega", "EVALUACIÓN", RED)
    add_bullet_list(slide, 0.5, 1.2, 9, 5.5, [
        "☐  Proyecto Supabase creado con los 4 esquemas (bronze, silver, gold, audit)",
        "☐  SQL ejecutado exitosamente — todas las tablas y vistas creadas",
        "☐  Datos Bronze con errores intencionales (cantidad NULL, vendedor vacío, ciudad sin normalizar)",
        "☐  Transformación Bronze → Silver funcionando (7 de 8 registros aprobados)",
        "☐  Vistas materializadas Gold con agregaciones de negocio",
        "☐  Catálogo de datos con metadatos de cada dataset",
        "☐  App Next.js desplegada en Vercel y funcionando",
        "☐  Las 11 tabs cargan datos correctamente",
        "☐  El dashboard SOLO consume datos de capa Gold",
        "☐  Poder explicar: ¿por qué se rechazó el registro #4?",
        "☐  Poder explicar: ¿qué hace INITCAP() y por qué es importante?",
        "☐  Poder explicar: ¿para qué sirve la auditoría de accesos?",
    ], font_size=13)

    # --- SLIDE 11: Cierre ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLATE_900)
    add_text_box(slide, 0.8, 2.5, 8.4, 1, "¡Manos a la obra!", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 3.8, 8.4, 1, "Repo: github.com/sjimenezlon/datalake-governance-hub\nApp: datalake-governance-hub.vercel.app", font_size=16, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 5.2, 8.4, 1, "Gestión y Gobernanza de Datos — EAFIT 2026", font_size=14, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    path = os.path.expanduser("~/Documents/datalake-governance-hub/PPTX_1_Paso_a_Paso_Actividad.pptx")
    prs.save(path)
    print(f"✅ PPTX 1 guardado: {path}")


# ================================================================
# PPTX 2: CÓMO FUNCIONA LA PLATAFORMA
# ================================================================
def create_como_funciona():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: Portada ---
    title_slide(prs,
        "Data Lake Governance Hub\nGuía de la Plataforma",
        "Cómo funciona cada sección con ejemplos reales\nSupabase + Next.js + Vercel")

    # --- SLIDE 2: Panorama general ---
    slide = content_slide(prs, "Panorama general — 11 secciones", "PLATAFORMA")
    add_bullet_list(slide, 0.5, 1.2, 4.3, 5.5, [
        "📊  Dashboard Gold — KPIs de negocio",
        "🔄  Bronze vs Silver — Antes y después",
        "✅  Calidad — Data Quality Dashboard",
        "🔗  Linaje — Trazabilidad end-to-end",
        "🥉  Bronze — Datos crudos (JSONB)",
        "🥈  Silver — Datos limpios y tipados",
    ], font_size=14)
    add_bullet_list(slide, 5, 1.2, 4.5, 5.5, [
        "📋  Catálogo — Metadatos obligatorios",
        "📡  Sensores IoT — Monitoreo en tiempo real",
        "🔒  Auditoría — Quién accedió a qué",
        "👥  Usuarios — RBAC por rol",
        "🏗️  Arquitectura — Stack técnico completo",
    ], font_size=14)
    add_text_box(slide, 0.5, 5.5, 9, 1, "URL: datalake-governance-hub.vercel.app", font_size=16, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # --- SLIDE 3: Dashboard Gold ---
    section_slide(prs, 1, "Dashboard Gold", "Solo consume datos de la capa Gold — Nunca toca Bronze ni Silver")

    slide = content_slide(prs, 'Tab: Dashboard Gold — "La vista del gerente"', "GOLD", GREEN)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "¿Qué muestra?", font_size=16, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 2, [
        "4 KPIs principales: Ingreso Total, Transacciones, Calidad Promedio, Registros Rechazados",
        "Tabla de Ventas por Ciudad con ingreso, ticket promedio y quality score",
        "Tabla de Ventas por Producto con unidades y precios",
    ], font_size=14)
    add_text_box(slide, 0.5, 3.8, 9, 0.5, "¿Por qué importa para Gobernanza?", font_size=16, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 4.3, 9, 2, [
        "El dashboard SOLO lee de gold_ventas_por_ciudad y gold_ventas_por_producto",
        "Nunca accede a Bronze ni Silver directamente — separación de capas",
        "Si Gold muestra un KPI incorrecto, el linaje permite rastrear hasta Bronze",
        "El quality score alerta si los datos tienen problemas de calidad",
    ], font_size=13)
    add_text_box(slide, 0.5, 6.2, 9, 0.5, "Ejemplo: Medellín tiene calidad 85% porque un registro tenía vendedor vacío → quality_score = 0.70", font_size=12, color=SLATE_400)

    # --- SLIDE 4: Bronze vs Silver ---
    section_slide(prs, 2, "Comparación Bronze vs Silver", "Ver exactamente qué se limpió y qué se rechazó")

    slide = content_slide(prs, "Tab: Bronze vs Silver — Antes y después", "TRANSFORMACIÓN", CYAN)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "¿Qué muestra?", font_size=16, color=CYAN, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 1.5, [
        "Cada registro Bronze lado a lado con su versión Silver",
        "Campos problemáticos resaltados en rojo (NULL, vacío, sin normalizar)",
        "Etiquetas: LIMPIO, TRANSFORMADO o RECHAZADO por registro",
    ], font_size=14)
    add_text_box(slide, 0.5, 3.5, 9, 0.5, "Ejemplos concretos que verán:", font_size=16, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 4.0, 9, 3, [
        'Registro #4: cantidad = NULL → RECHAZADO (no pasa a Silver)',
        'Registro #4: vendedor = "" → RECHAZADO (mismo registro)',
        'Registro #4: ciudad = "medellin" → RECHAZADO (mismo registro)',
        'Registro #7: vendedor = "Juan Pérez" → LIMPIO, quality_score = 1.00',
        'En otro registro: ciudad "medellin" se normalizaría → "Medellin" con INITCAP()',
        'Vendedor vacío en otros registros → "Sin asignar" con quality_score = 0.70',
    ], font_size=13)

    # --- SLIDE 5: Calidad ---
    section_slide(prs, 3, "Data Quality Dashboard", "Las 6 dimensiones de calidad de DAMA-DMBOK2")

    slide = content_slide(prs, "Tab: Calidad — Monitoreo automatizado", "CALIDAD", AMBER)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Métricas que muestra:", font_size=16, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 4.3, 2.5, [
        "Registros Bronze: 8",
        "Aprobados Silver: 7",
        "Rechazados: 1",
        "Tasa de aprobación: 87%",
        "Quality Score promedio: 96%",
    ], font_size=13)
    add_text_box(slide, 5, 1.2, 4.5, 0.5, "7 reglas de calidad activas:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 5, 1.8, 4.5, 2.5, [
        "Completitud: cantidad NOT NULL",
        "Validez: cantidad > 0, precio > 0",
        "Consistencia: INITCAP(ciudad)",
        "Unicidad: deduplicación por ID",
        "Confiabilidad: quality_score < 0.8 = alerta",
    ], font_size=12)
    add_text_box(slide, 0.5, 4.5, 9, 0.5, "Problemas detectados automáticamente en Bronze:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 5.0, 9, 2, [
        '🔴 CRÍTICO: cantidad es NULL en registro #4',
        '🟡 WARNING: vendedor vacío en registro #4',
        '🔵 INFO: ciudad "medellin" no normalizada en registro #4',
    ], font_size=13)
    add_text_box(slide, 0.5, 6.5, 9, 0.5, "Barra visual: Quality Score por dataset del catálogo (80% a 98%)", font_size=12, color=SLATE_400)

    # --- SLIDE 6: Linaje ---
    section_slide(prs, 4, "Linaje de Datos (Data Lineage)", "Trazabilidad completa: ¿de dónde viene cada dato?")

    slide = content_slide(prs, "Tab: Linaje — Grafo y trazabilidad por registro", "LINAJE", PURPLE)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Grafo de linaje (flujo completo):", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 1.5, [
        "FUENTE (SAP ERP, IoT Gateway) → BRONZE (JSONB inmutable) → SILVER (tipado, validado) → GOLD (agregado) → CONSUMO (Dashboard, Analistas, ML)",
        "Cada nivel muestra las tablas específicas y conteo de registros",
    ], font_size=13)
    add_text_box(slide, 0.5, 3.5, 9, 0.5, "Linaje a nivel de registro:", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 4.0, 9, 2.5, [
        "Cada registro Silver tiene un source_record_id que apunta a su registro Bronze",
        "Ejemplo: Silver #1 (Sensor IoT X200, Medellín) ← Bronze #1 (SAP-ERP, export_ventas_20260301.json)",
        "Si un KPI de Gold parece incorrecto: Gold → Silver (qué registros lo componen) → Bronze (dato original)",
        "En producción se usa Apache Atlas, OpenLineage o dbt lineage para esto automáticamente",
    ], font_size=13)

    # --- SLIDE 7: Catálogo ---
    section_slide(prs, 5, "Catálogo de Datos", "Sin catálogo = Data Swamp")

    slide = content_slide(prs, "Tab: Catálogo — Metadatos obligatorios", "CATÁLOGO", BLUE)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Cada dataset registra:", font_size=14, color=BLUE, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 4.3, 3, [
        "Esquema y nombre de tabla",
        "Descripción en lenguaje natural",
        "Dueño (owner) responsable",
        "Clasificación: public, internal, confidential",
    ], font_size=13)
    add_bullet_list(slide, 5, 1.8, 4.5, 3, [
        "Indicador de PII (datos personales)",
        "Conteo de filas",
        "Freshness (frecuencia de actualización)",
        "Quality Score con barra visual",
    ], font_size=13)
    add_text_box(slide, 0.5, 4.5, 9, 0.5, "Ejemplo concreto:", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 5.0, 9, 2, [
        "bronze.ventas_raw → Confidential, PII=Sí, Owner=equipo-datos, Quality=85%, Freshness=daily",
        "gold.ventas_por_ciudad → Internal, PII=No, Owner=equipo-bi, Quality=98%, Freshness=daily",
        "gold.estado_sensores → Public, PII=No, Owner=equipo-iot, Quality=97%, Freshness=hourly",
    ], font_size=12)

    # --- SLIDE 8: Auditoría ---
    section_slide(prs, 6, "Auditoría de Accesos", "Requisito de cumplimiento: Habeas Data, GDPR, SOX")

    slide = content_slide(prs, "Tab: Auditoría — Trazabilidad de accesos", "AUDITORÍA", RED)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Cada acceso registra:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 1.5, [
        "Timestamp exacto, usuario (email), rol, acción (SELECT/INSERT/EXPORT), tabla, número de filas, IP",
    ], font_size=13)
    add_text_box(slide, 0.5, 3.5, 9, 0.5, "Ejemplos reales en la plataforma:", font_size=14, color=AMBER, bold=True)
    add_bullet_list(slide, 0.5, 4.0, 9, 3, [
        "carlos.ingeniero (data_engineer) → INSERT en bronze.ventas_raw → 8 filas → IP: 10.0.1.15",
        "carlos.ingeniero (data_engineer) → INSERT en silver.ventas_clean → 7 filas → IP: 10.0.1.15",
        "maria.ciencia (data_scientist) → SELECT en silver.ventas_clean → 7 filas → IP: 10.0.2.22",
        "pedro.analista (analyst) → SELECT en gold.ventas_por_ciudad → 5 filas → IP: 10.0.3.8",
        "pedro.analista (analyst) → EXPORT en gold.ventas_por_ciudad → 5 filas → IP: 10.0.3.8",
    ], font_size=12)
    add_text_box(slide, 0.5, 6.3, 9, 0.5, "Nota: El analista SOLO accede a Gold. El data engineer accede a Bronze y Silver. Separación de capas.", font_size=12, color=SLATE_400)

    # --- SLIDE 9: Usuarios y RBAC ---
    section_slide(prs, 7, "Usuarios y Control de Acceso (RBAC)", "Principio de mínimo privilegio")

    slide = content_slide(prs, "Tab: Usuarios — 4 roles con permisos diferenciados", "RBAC", CYAN)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Roles y permisos:", font_size=14, color=CYAN, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 3.5, [
        "ADMIN → Bronze + Silver + Gold + Audit + Admin — control total",
        "DATA ENGINEER → Bronze (R/W) + Silver (R/W) + Gold (R) — construye pipelines",
        "DATA SCIENTIST → Silver masked (R) + Gold (R) + Sandbox (R/W) — explora sin ver PII",
        "ANALYST → Gold (R) — solo lectura de la capa de negocio, datos ya agregados",
    ], font_size=14)
    add_text_box(slide, 0.5, 4.5, 9, 0.5, "Seguridad implementada:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 5.0, 9, 2, [
        "Row Level Security (RLS) en tablas Silver — PostgreSQL nativo",
        "Vista ventas_masked: NIT \"900123456\" → \"900****56\", nombre enmascarado",
        "Data Scientists ven ventas_masked, NO ventas_clean — no acceden al NIT real",
        "Cifrado: Supabase aplica TLS (en tránsito) y AES-256 (en reposo) por defecto",
    ], font_size=13)

    # --- SLIDE 10: Sensores IoT ---
    slide = content_slide(prs, "Tab: Sensores IoT — Monitoreo desde Gold", "IOT", GREEN)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Cada sensor muestra:", font_size=14, color=GREEN, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 2, [
        "ID del sensor, ubicación, tipo de medición (caudal, presión, calidad, nivel)",
        "Valor promedio, mínimo y máximo, unidad de medida",
        "Nivel de batería con código de colores (verde >70%, ámbar >50%, rojo <50%)",
        "Indicador de anomalía — SEN-001 tiene valor -5.0 L/s → anomalía detectada",
    ], font_size=13)
    add_text_box(slide, 0.5, 4.0, 9, 0.5, "Ejemplo de anomalía:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 4.5, 9, 2, [
        "SEN-001 (Planta Norte) midió caudal = -5.0 L/s — un caudal negativo es físicamente imposible",
        "El pipeline Bronze→Silver marcó is_anomaly = TRUE para esa lectura",
        "En Gold, el campo tiene_anomalias = TRUE para SEN-001, alertando al operador",
        "SEN-004 (Tanque Reserva) tiene batería al 45% — alerta de mantenimiento",
    ], font_size=13)

    # --- SLIDE 11: Arquitectura ---
    slide = content_slide(prs, "Tab: Arquitectura — Vista técnica completa", "ARQUITECTURA", PURPLE)
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Flujo Medallion visual:", font_size=14, color=PURPLE, bold=True)
    add_bullet_list(slide, 0.5, 1.8, 9, 1.5, [
        "Bronze (8 reg, JSONB) → Silver (7 reg, tipado) → Gold (agregaciones, vistas materializadas)",
        "Cada zona muestra sus tablas y conteo de registros en tiempo real",
    ], font_size=13)
    add_text_box(slide, 0.5, 3.5, 4.3, 0.5, "Seguridad implementada:", font_size=14, color=RED, bold=True)
    add_bullet_list(slide, 0.5, 4.0, 4.3, 2.5, [
        "✓ Row Level Security (RLS)",
        "✓ Enmascaramiento PII",
        "✓ RBAC — 4 roles",
        "✓ Auditoría con IP",
        "✓ Cifrado TLS + AES-256",
    ], font_size=12)
    add_text_box(slide, 5, 3.5, 4.5, 0.5, "Gobernanza implementada:", font_size=14, color=BLUE, bold=True)
    add_bullet_list(slide, 5, 4.0, 4.5, 2.5, [
        "✓ Catálogo con metadatos",
        "✓ Clasificación de datos",
        "✓ Quality Score por dataset",
        "✓ Linaje Bronze→Silver→Gold",
        "✓ Data Stewardship (owners)",
    ], font_size=12)
    add_text_box(slide, 0.5, 6.2, 9, 0.5, "Stack: Supabase (PostgreSQL+RLS) + Next.js 14 + Tailwind CSS + Vercel (Edge)", font_size=13, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    # --- SLIDE 12: Cierre ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLATE_900)
    add_text_box(slide, 0.8, 2.0, 8.4, 1, "Explora la plataforma", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 3.3, 8.4, 1.5, "datalake-governance-hub.vercel.app\n\nNavega por las 11 tabs, compara Bronze vs Silver,\nrevisa la auditoría, explora el catálogo.", font_size=16, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 0.8, 5.5, 8.4, 1, "Gestión y Gobernanza de Datos — EAFIT 2026", font_size=14, color=SLATE_400, alignment=PP_ALIGN.CENTER)

    path = os.path.expanduser("~/Documents/datalake-governance-hub/PPTX_2_Como_Funciona_Plataforma.pptx")
    prs.save(path)
    print(f"✅ PPTX 2 guardado: {path}")


# ================================================================
if __name__ == '__main__':
    create_paso_a_paso()
    create_como_funciona()
    print("\n🎉 Ambos PowerPoints creados exitosamente!")
